#!/usr/bin/env python3
"""
Full-Featured Markdown to PowerPoint Converter
Supports advanced features, multiple themes, and template-based conversion with tag replacement
"""

import re
import yaml
import json
import argparse
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass, field
from enum import Enum
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False


# ============================================================================
# THEMES CONFIGURATION
# ============================================================================

class Theme:
    """Theme configuration for presentations"""
    
    def __init__(self, name: str, config: Dict):
        self.name = name
        self.background = config.get("background", "FFFFFF")
        self.title_color = config.get("title_color", "000000")
        self.text_color = config.get("text_color", "333333")
        self.accent_color = config.get("accent_color", "0066CC")
        self.title_font = config.get("title_font", "Calibri")
        self.body_font = config.get("body_font", "Calibri")
        self.title_size = config.get("title_size", 44)
        self.body_size = config.get("body_size", 20)
        self.bullet_style = config.get("bullet_style", "•")


THEMES = {
    "default": Theme("Default", {
        "background": "FFFFFF",
        "title_color": "1F4E78",
        "text_color": "2F2F2F",
        "accent_color": "4472C4",
        "title_font": "Calibri Light",
        "body_font": "Calibri",
        "title_size": 44,
        "body_size": 20,
        "bullet_style": "•"
    }),
    
    "corporate": Theme("Corporate", {
        "background": "FFFFFF",
        "title_color": "003366",
        "text_color": "333333",
        "accent_color": "0066CC",
        "title_font": "Arial Black",
        "body_font": "Arial",
        "title_size": 40,
        "body_size": 18,
        "bullet_style": "■"
    }),
    
    "modern": Theme("Modern", {
        "background": "F8F9FA",
        "title_color": "2C3E50",
        "text_color": "34495E",
        "accent_color": "3498DB",
        "title_font": "Segoe UI",
        "body_font": "Segoe UI",
        "title_size": 48,
        "body_size": 22,
        "bullet_style": "▸"
    }),
    
    "minimal": Theme("Minimal", {
        "background": "FFFFFF",
        "title_color": "000000",
        "text_color": "404040",
        "accent_color": "808080",
        "title_font": "Helvetica",
        "body_font": "Helvetica",
        "title_size": 54,
        "body_size": 24,
        "bullet_style": "–"
    }),
    
    "creative": Theme("Creative", {
        "background": "FFF9E6",
        "title_color": "E74C3C",
        "text_color": "2C3E50",
        "accent_color": "9B59B6",
        "title_font": "Georgia",
        "body_font": "Georgia",
        "title_size": 46,
        "body_size": 20,
        "bullet_style": "★"
    }),
    
    "dark": Theme("Dark", {
        "background": "1E1E1E",
        "title_color": "FFFFFF",
        "text_color": "E0E0E0",
        "accent_color": "00D9FF",
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": 44,
        "body_size": 20,
        "bullet_style": "›"
    }),
    
    "academic": Theme("Academic", {
        "background": "FFFFFF",
        "title_color": "1A237E",
        "text_color": "37474F",
        "accent_color": "3F51B5",
        "title_font": "Times New Roman",
        "body_font": "Times New Roman",
        "title_size": 40,
        "body_size": 18,
        "bullet_style": "•"
    }),
    
    "tech": Theme("Tech", {
        "background": "0A0E27",
        "title_color": "00FF88",
        "text_color": "FFFFFF",
        "accent_color": "FF006E",
        "title_font": "Consolas",
        "body_font": "Consolas",
        "title_size": 42,
        "body_size": 18,
        "bullet_style": ">"
    }),
    
    "nature": Theme("Nature", {
        "background": "F1F8F4",
        "title_color": "2D5016",
        "text_color": "3E5C3E",
        "accent_color": "6AA84F",
        "title_font": "Trebuchet MS",
        "body_font": "Trebuchet MS",
        "title_size": 44,
        "body_size": 20,
        "bullet_style": "🌿"
    }),
    
    "sunset": Theme("Sunset", {
        "background": "FFF4E6",
        "title_color": "D84315",
        "text_color": "4E342E",
        "accent_color": "FF6F00",
        "title_font": "Palatino Linotype",
        "body_font": "Palatino Linotype",
        "title_size": 46,
        "body_size": 20,
        "bullet_style": "●"
    })
}


# ============================================================================
# SLIDE TYPES
# ============================================================================

class SlideType(Enum):
    """Supported slide types"""
    TITLE = "title"
    TITLE_AND_CONTENT = "title_content"
    TWO_COLUMN = "two_column"
    SECTION = "section"
    BLANK = "blank"
    IMAGE_FULL = "image_full"
    IMAGE_LEFT = "image_left"
    IMAGE_RIGHT = "image_right"
    QUOTE = "quote"
    TABLE = "table"
    CHART = "chart"
    COMPARISON = "comparison"
    TIMELINE = "timeline"
    CLOSING = "closing"


# ============================================================================
# DATA CLASSES
# ============================================================================

@dataclass
class SlideConfig:
    """Configuration for a single slide"""
    type: SlideType
    title: Optional[str] = None
    subtitle: Optional[str] = None
    content: Optional[List[str]] = None
    left_content: Optional[List[str]] = None
    right_content: Optional[List[str]] = None
    image_path: Optional[str] = None
    table_data: Optional[List[List[str]]] = None
    chart_data: Optional[Dict] = None
    quote_text: Optional[str] = None
    quote_author: Optional[str] = None
    notes: Optional[str] = None
    background_color: Optional[str] = None
    metadata: Dict = field(default_factory=dict)


@dataclass
class PresentationConfig:
    """Global presentation configuration"""
    title: str = "Presentation"
    author: Optional[str] = None
    theme: str = "default"
    aspect_ratio: str = "16:9"
    slide_numbers: bool = True
    company: Optional[str] = None
    logo_path: Optional[str] = None
    footer_text: Optional[str] = None
    metadata: Dict = field(default_factory=dict)


# ============================================================================
# MARKDOWN PARSER
# ============================================================================

class MarkdownParser:
    """Parse markdown content into presentation structure"""
    
    PATTERNS = {
        "frontmatter": r"^---\s*\n(.*?)\n---\s*\n",
        "slide_break": r"\n---+\n",
        "heading1": r"^# (.+)$",
        "heading2": r"^## (.+)$",
        "heading3": r"^### (.+)$",
        "bullet": r"^[\*\-\+] (.+)$",
        "numbered": r"^\d+\. (.+)$",
        "image": r"!\[(.*?)\]\((.*?)\)",
        "link": r"\[(.*?)\]\((.*?)\)",
        "bold": r"\*\*(.+?)\*\*",
        "italic": r"\*(.+?)\*",
        "code": r"`(.+?)`",
        "table_row": r"^\|(.+)\|$",
        "quote": r"^> (.+)$",
        "slide_type": r"<!-- slide: (\w+) -->"
    }
    
    def __init__(self):
        self.config = PresentationConfig()
        self.slides = []
    
    def parse(self, markdown_content: str) -> Tuple[PresentationConfig, List[SlideConfig]]:
        """Parse markdown content"""
        # Extract frontmatter
        self._parse_frontmatter(markdown_content)
        
        # Remove frontmatter from content
        content = re.sub(self.PATTERNS["frontmatter"], "", markdown_content, flags=re.MULTILINE | re.DOTALL)
        
        # Split into slides
        slide_texts = re.split(self.PATTERNS["slide_break"], content)
        
        for slide_text in slide_texts:
            if slide_text.strip():
                slide = self._parse_slide(slide_text.strip())
                if slide:
                    self.slides.append(slide)
        
        return self.config, self.slides
    
    def _parse_frontmatter(self, content: str):
        """Parse YAML frontmatter"""
        match = re.search(self.PATTERNS["frontmatter"], content, re.MULTILINE | re.DOTALL)
        if match:
            try:
                metadata = yaml.safe_load(match.group(1))
                self.config = PresentationConfig(
                    title=metadata.get("title", "Presentation"),
                    author=metadata.get("author"),
                    theme=metadata.get("theme", "default"),
                    aspect_ratio=metadata.get("aspect_ratio", "16:9"),
                    slide_numbers=metadata.get("slide_numbers", True),
                    company=metadata.get("company"),
                    logo_path=metadata.get("logo_path"),
                    footer_text=metadata.get("footer_text"),
                    metadata=metadata
                )
            except yaml.YAMLError:
                pass
    
    def _parse_slide(self, slide_text: str) -> Optional[SlideConfig]:
        """Parse a single slide"""
        lines = slide_text.split("\n")
        
        # Detect slide type from comment
        slide_type = SlideType.TITLE_AND_CONTENT
        type_match = re.search(self.PATTERNS["slide_type"], slide_text)
        if type_match:
            try:
                slide_type = SlideType(type_match.group(1))
            except ValueError:
                pass
        
        title = None
        subtitle = None
        content = []
        left_content = []
        right_content = []
        image_path = None
        quote_text = None
        quote_author = None
        table_data = []
        in_left_column = False
        in_right_column = False
        
        for line in lines:
            line = line.strip()
            
            # Skip empty lines and comments
            if not line or line.startswith("<!--"):
                continue
            
            # Column markers
            if line == "::left::":
                in_left_column = True
                in_right_column = False
                slide_type = SlideType.TWO_COLUMN
                continue
            elif line == "::right::":
                in_left_column = False
                in_right_column = True
                slide_type = SlideType.TWO_COLUMN
                continue
            
            # Headings
            h1_match = re.match(self.PATTERNS["heading1"], line)
            if h1_match:
                if not title:
                    title = h1_match.group(1)
                else:
                    content.append(h1_match.group(1))
                continue
            
            h2_match = re.match(self.PATTERNS["heading2"], line)
            if h2_match:
                if not title:
                    title = h2_match.group(1)
                elif not subtitle:
                    subtitle = h2_match.group(1)
                else:
                    content.append(h2_match.group(1))
                continue
            
            h3_match = re.match(self.PATTERNS["heading3"], line)
            if h3_match:
                if not subtitle and title:
                    subtitle = h3_match.group(1)
                else:
                    content.append(h3_match.group(1))
                continue
            
            # Quotes
            quote_match = re.match(self.PATTERNS["quote"], line)
            if quote_match:
                if not quote_text:
                    quote_text = quote_match.group(1)
                    slide_type = SlideType.QUOTE
                else:
                    quote_author = quote_match.group(1).lstrip("— -")
                continue
            
            # Images
            image_match = re.search(self.PATTERNS["image"], line)
            if image_match:
                image_path = image_match.group(2)
                if "![full]" in line or "![fullscreen]" in line:
                    slide_type = SlideType.IMAGE_FULL
                continue
            
            # Bullets
            bullet_match = re.match(self.PATTERNS["bullet"], line)
            if bullet_match:
                item = bullet_match.group(1)
                if in_left_column:
                    left_content.append(item)
                elif in_right_column:
                    right_content.append(item)
                else:
                    content.append(item)
                continue
            
            # Numbered lists
            numbered_match = re.match(self.PATTERNS["numbered"], line)
            if numbered_match:
                item = numbered_match.group(1)
                if in_left_column:
                    left_content.append(item)
                elif in_right_column:
                    right_content.append(item)
                else:
                    content.append(item)
                continue
            
            # Table rows
            if line.startswith("|"):
                cells = [cell.strip() for cell in line.strip("|").split("|")]
                if cells and not all(c.replace("-", "").strip() == "" for c in cells):
                    table_data.append(cells)
                    slide_type = SlideType.TABLE
                continue
            
            # Regular text
            if line:
                if in_left_column:
                    left_content.append(line)
                elif in_right_column:
                    right_content.append(line)
                else:
                    content.append(line)
        
        # Determine slide type if not explicitly set
        if not title and not content:
            return None
        
        if quote_text:
            slide_type = SlideType.QUOTE
        elif table_data:
            slide_type = SlideType.TABLE
        elif not title and not content:
            slide_type = SlideType.BLANK
        elif title and subtitle and not content:
            slide_type = SlideType.TITLE
        
        return SlideConfig(
            type=slide_type,
            title=title,
            subtitle=subtitle,
            content=content if content else None,
            left_content=left_content if left_content else None,
            right_content=right_content if right_content else None,
            image_path=image_path,
            table_data=table_data if table_data else None,
            quote_text=quote_text,
            quote_author=quote_author
        )


# ============================================================================
# POWERPOINT GENERATOR
# ============================================================================

class PowerPointGenerator:
    """Generate PowerPoint presentations from parsed data"""
    
    def __init__(self, theme_name: str = "default"):
        self.theme = THEMES.get(theme_name, THEMES["default"])
        self.prs = None
    
    def generate(self, config: PresentationConfig, slides: List[SlideConfig]) -> Presentation:
        """Generate PowerPoint presentation"""
        # Create presentation
        self.prs = Presentation()
        
        # Set slide dimensions
        if config.aspect_ratio == "16:9":
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        else:  # 4:3
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
        
        # Generate slides
        for slide_config in slides:
            self._create_slide(slide_config)
        
        # Add footer if specified
        if config.footer_text or config.slide_numbers:
            self._add_footer(config.footer_text, config.slide_numbers)
        
        return self.prs
    
    def _create_slide(self, config: SlideConfig):
        """Create a single slide based on type"""
        if config.type == SlideType.TITLE:
            self._create_title_slide(config)
        elif config.type == SlideType.SECTION:
            self._create_section_slide(config)
        elif config.type == SlideType.TWO_COLUMN:
            self._create_two_column_slide(config)
        elif config.type == SlideType.QUOTE:
            self._create_quote_slide(config)
        elif config.type == SlideType.TABLE:
            self._create_table_slide(config)
        elif config.type == SlideType.IMAGE_FULL:
            self._create_image_slide(config, full=True)
        elif config.type == SlideType.BLANK:
            self._create_blank_slide(config)
        else:
            self._create_content_slide(config)
    
    def _create_title_slide(self, config: SlideConfig):
        """Create title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Set background
        self._set_background(slide, config.background_color)
        
        # Add title
        if config.title:
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(11.33), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = config.title
            title_frame.paragraphs[0].font.size = Pt(self.theme.title_size)
            title_frame.paragraphs[0].font.name = self.theme.title_font
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.title_color)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        if config.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(11.33), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = config.subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(self.theme.body_size + 4)
            subtitle_frame.paragraphs[0].font.name = self.theme.body_font
            subtitle_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.text_color)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_section_slide(self, config: SlideConfig):
        """Create section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Set background with accent color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(self.theme.accent_color)
        
        # Add section title
        if config.title:
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(11.33), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = config.title
            title_frame.paragraphs[0].font.size = Pt(self.theme.title_size + 8)
            title_frame.paragraphs[0].font.name = self.theme.title_font
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, config: SlideConfig):
        """Create content slide with bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Set background
        self._set_background(slide, config.background_color)
        
        # Add title
        if config.title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = config.title
            title_frame.paragraphs[0].font.size = Pt(self.theme.title_size - 6)
            title_frame.paragraphs[0].font.name = self.theme.title_font
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.title_color)
        
        # Add content
        if config.content:
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.5)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, item in enumerate(config.content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"{self.theme.bullet_style} {item}"
                p.font.size = Pt(self.theme.body_size)
                p.font.name = self.theme.body_font
                p.font.color.rgb = self._hex_to_rgb(self.theme.text_color)
                p.level = 0
                p.space_before = Pt(6)
    
    def _create_two_column_slide(self, config: SlideConfig):
        """Create two-column slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Set background
        self._set_background(slide, config.background_color)
        
        # Add title
        if config.title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = config.title
            title_frame.paragraphs[0].font.size = Pt(self.theme.title_size - 6)
            title_frame.paragraphs[0].font.name = self.theme.title_font
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.title_color)
        
        # Left column
        if config.left_content:
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(6), Inches(5.5)
            )
            left_frame = left_box.text_frame
            for i, item in enumerate(config.left_content):
                if i == 0:
                    p = left_frame.paragraphs[0]
                else:
                    p = left_frame.add_paragraph()
                p.text = f"{self.theme.bullet_style} {item}"
                p.font.size = Pt(self.theme.body_size - 2)
                p.font.name = self.theme.body_font
                p.font.color.rgb = self._hex_to_rgb(self.theme.text_color)
        
        # Right column
        if config.right_content:
            right_box = slide.shapes.add_textbox(
                Inches(6.83), Inches(1.5), Inches(6), Inches(5.5)
            )
            right_frame = right_box.text_frame
            for i, item in enumerate(config.right_content):
                if i == 0:
                    p = right_frame.paragraphs[0]
                else:
                    p = right_frame.add_paragraph()
                p.text = f"{self.theme.bullet_style} {item}"
                p.font.size = Pt(self.theme.body_size - 2)
                p.font.name = self.theme.body_font
                p.font.color.rgb = self._hex_to_rgb(self.theme.text_color)
    
    def _create_quote_slide(self, config: SlideConfig):
        """Create quote slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Set background
        self._set_background(slide, config.background_color)
        
        # Add quote
        if config.quote_text:
            quote_box = slide.shapes.add_textbox(
                Inches(2), Inches(2.5), Inches(9.33), Inches(2)
            )
            quote_frame = quote_box.text_frame
            quote_frame.text = f'"{config.quote_text}"'
            quote_frame.paragraphs[0].font.size = Pt(self.theme.body_size + 8)
            quote_frame.paragraphs[0].font.name = self.theme.body_font
            quote_frame.paragraphs[0].font.italic = True
            quote_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.text_color)
            quote_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add author
        if config.quote_author:
            author_box = slide.shapes.add_textbox(
                Inches(2), Inches(5), Inches(9.33), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_frame.text = f"— {config.quote_author}"
            author_frame.paragraphs[0].font.size = Pt(self.theme.body_size)
            author_frame.paragraphs[0].font.name = self.theme.body_font
            author_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.accent_color)
            author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_table_slide(self, config: SlideConfig):
        """Create table slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Set background
        self._set_background(slide, config.background_color)
        
        # Add title
        if config.title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = config.title
            title_frame.paragraphs[0].font.size = Pt(self.theme.title_size - 6)
            title_frame.paragraphs[0].font.name = self.theme.title_font
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.title_color)
        
        # Add table
        if config.table_data and len(config.table_data) > 0:
            rows = len(config.table_data)
            cols = len(config.table_data[0])
            
            left = Inches(1)
            top = Inches(1.8)
            width = Inches(11.33)
            height = Inches(5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Fill table
            for i, row_data in enumerate(config.table_data):
                for j, cell_data in enumerate(row_data):
                    if j < len(table.columns):
                        cell = table.cell(i, j)
                        cell.text = cell_data
                        
                        # Header row styling
                        if i == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = self._hex_to_rgb(self.theme.accent_color)
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.bold = True
                                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                                paragraph.font.size = Pt(self.theme.body_size - 2)
                        else:
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(self.theme.body_size - 4)
                                paragraph.font.color.rgb = self._hex_to_rgb(self.theme.text_color)
    
    def _create_image_slide(self, config: SlideConfig, full: bool = False):
        """Create image slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Set background
        self._set_background(slide, config.background_color)
        
        # Add image
        if config.image_path:
            try:
                if full:
                    # Full screen image
                    slide.shapes.add_picture(
                        config.image_path,
                        Inches(0), Inches(0),
                        width=self.prs.slide_width,
                        height=self.prs.slide_height
                    )
                else:
                    # Centered image with title
                    if config.title:
                        title_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8)
                        )
                        title_frame = title_box.text_frame
                        title_frame.text = config.title
                        title_frame.paragraphs[0].font.size = Pt(self.theme.title_size - 6)
                        title_frame.paragraphs[0].font.name = self.theme.title_font
                        title_frame.paragraphs[0].font.bold = True
                        title_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.title_color)
                    
                    slide.shapes.add_picture(
                        config.image_path,
                        Inches(2), Inches(1.5),
                        width=Inches(9.33), height=Inches(5.5)
                    )
            except Exception as e:
                print(f"Warning: Could not add image {config.image_path}: {e}")
    
    def _create_blank_slide(self, config: SlideConfig):
        """Create blank slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, config.background_color)
    
    def _set_background(self, slide, color: Optional[str] = None):
        """Set slide background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        if color:
            fill.fore_color.rgb = self._hex_to_rgb(color)
        else:
            fill.fore_color.rgb = self._hex_to_rgb(self.theme.background)
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGB"""
        hex_color = hex_color.lstrip("#")
        return RGBColor(*[int(hex_color[i:i+2], 16) for i in (0, 2, 4)])
    
    def _add_footer(self, text: Optional[str], slide_numbers: bool):
        """Add footer to all slides"""
        for i, slide in enumerate(self.prs.slides):
            footer_text = []
            if text:
                footer_text.append(text)
            if slide_numbers:
                footer_text.append(f"{i + 1}")
            
            if footer_text:
                footer_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(7), Inches(12.33), Inches(0.3)
                )
                footer_frame = footer_box.text_frame
                footer_frame.text = " | ".join(footer_text)
                footer_frame.paragraphs[0].font.size = Pt(10)
                footer_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.theme.text_color)
                footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ============================================================================
# TEMPLATE PROCESSOR
# ============================================================================

class TemplateProcessor:
    """Process PowerPoint templates with tag replacement"""
    
    def __init__(self, template_path: str):
        self.template_path = template_path
        self.prs = None
    
    def process(self, replacements: Dict[str, str]) -> Presentation:
        """Process template with tag replacements"""
        # Load template
        self.prs = Presentation(self.template_path)
        
        # Process all slides
        for slide in self.prs.slides:
            self._process_slide(slide, replacements)
        
        return self.prs
    
    def _process_slide(self, slide, replacements: Dict[str, str]):
        """Process a single slide"""
        # Process all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                self._replace_text_in_frame(shape.text_frame, replacements)
            
            # Process tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        self._replace_text_in_frame(cell.text_frame, replacements)
        
        # Process notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            self._replace_text_in_frame(notes_slide.notes_text_frame, replacements)
    
    def _replace_text_in_frame(self, text_frame, replacements: Dict[str, str]):
        """Replace tags in text frame"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                modified_text = original_text
                
                # Replace all tags
                for tag, value in replacements.items():
                    # Support multiple tag formats: {{tag}}, {tag}, [tag]
                    patterns = [
                        f"{{{{{tag}}}}}",  # {{tag}}
                        f"{{{tag}}}",      # {tag}
                        f"[{tag}]",        # [tag]
                        f"${tag}",         # $tag
                        f"%{tag}%"         # %tag%
                    ]
                    
                    for pattern in patterns:
                        modified_text = modified_text.replace(pattern, str(value))
                
                run.text = modified_text


# ============================================================================
# MAIN CONVERTER
# ============================================================================

class MarkdownToPowerPoint:
    """Main converter class"""
    
    def __init__(self, theme: str = "default"):
        self.parser = MarkdownParser()
        self.theme = theme
    
    def convert(self, markdown_content: str, output_path: str):
        """Convert markdown to PowerPoint"""
        # Parse markdown
        config, slides = self.parser.parse(markdown_content)
        
        # Use theme from config if specified
        theme = config.theme if config.theme in THEMES else self.theme
        
        # Generate presentation
        generator = PowerPointGenerator(theme)
        prs = generator.generate(config, slides)
        
        # Save
        prs.save(output_path)
        
        print(f"✅ Presentation created: {output_path}")
        print(f"   Theme: {theme}")
        print(f"   Slides: {len(prs.slides)}")
        print(f"   Aspect Ratio: {config.aspect_ratio}")
    
    def convert_file(self, markdown_file: str, output_path: str):
        """Convert markdown file to PowerPoint"""
        with open(markdown_file, 'r', encoding='utf-8') as f:
            markdown_content = f.read()
        self.convert(markdown_content, output_path)
    
    def process_template(self, template_path: str, replacements: Dict[str, str], output_path: str):
        """Process template with tag replacements"""
        processor = TemplateProcessor(template_path)
        prs = processor.process(replacements)
        prs.save(output_path)
        
        print(f"✅ Template processed: {output_path}")
        print(f"   Replacements: {len(replacements)} tags")
    
    def process_template_from_json(self, template_path: str, json_file: str, output_path: str):
        """Process template with replacements from JSON file"""
        with open(json_file, 'r', encoding='utf-8') as f:
            replacements = json.load(f)
        self.process_template(template_path, replacements, output_path)


# ============================================================================
# CLI
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Convert Markdown to PowerPoint with themes and template support"
    )
    
    subparsers = parser.add_subparsers(dest="command", help="Command to execute")
    
    # Convert command
    convert_parser = subparsers.add_parser("convert", help="Convert markdown to PowerPoint")
    convert_parser.add_argument("input", help="Input markdown file")
    convert_parser.add_argument("output", help="Output PowerPoint file")
    convert_parser.add_argument("--theme", default="default", choices=list(THEMES.keys()),
                                help="Theme to use")
    
    # Template command
    template_parser = subparsers.add_parser("template", help="Process template with tag replacement")
    template_parser.add_argument("template", help="Template PowerPoint file")
    template_parser.add_argument("output", help="Output PowerPoint file")
    template_parser.add_argument("--json", help="JSON file with replacements")
    template_parser.add_argument("--tag", action="append", nargs=2, metavar=("KEY", "VALUE"),
                                 help="Tag replacement (can be used multiple times)")
    
    # List themes command
    list_parser = subparsers.add_parser("themes", help="List available themes")
    
    args = parser.parse_args()
    
    if args.command == "convert":
        converter = MarkdownToPowerPoint(theme=args.theme)
        converter.convert_file(args.input, args.output)
    
    elif args.command == "template":
        converter = MarkdownToPowerPoint()
        
        if args.json:
            converter.process_template_from_json(args.template, args.json, args.output)
        elif args.tag:
            replacements = {key: value for key, value in args.tag}
            converter.process_template(args.template, replacements, args.output)
        else:
            print("Error: Must provide either --json or --tag arguments")
            return 1
    
    elif args.command == "themes":
        print("\n📚 Available Themes:\n")
        for name, theme in THEMES.items():
            print(f"  • {name:12s} - {theme.name}")
            print(f"    Colors: Title={theme.title_color}, Text={theme.text_color}, Accent={theme.accent_color}")
            print(f"    Fonts: {theme.title_font} / {theme.body_font}")
            print()
    
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
