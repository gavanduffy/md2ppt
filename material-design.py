#!/usr/bin/env python3
"""
Material Design Theme Extension for PowerPoint MCP Server
"""

import asyncio
import json
from typing import Dict, List, Tuple, Optional, Any
from dataclasses import dataclass
from enum import Enum
import colorsys
import random

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Try to import MCP types, but allow fallback for standalone use
try:
    from mcp import types
except ImportError:
    # Create minimal type stubs for standalone use
    class types:
        class Tool:
            def __init__(self, name, description, inputSchema):
                self.name = name
                self.description = description
                self.inputSchema = inputSchema
        
        class TextContent:
            def __init__(self, type, text):
                self.type = type
                self.text = text

class MaterialColorPalette(Enum):
    """Material Design 3 Color Palettes"""

    # Primary Colors
    RED = {
        "50": "FFEBEE", "100": "FFCDD2", "200": "EF9A9A", "300": "E57373",
        "400": "EF5350", "500": "F44336", "600": "E53935", "700": "D32F2F",
        "800": "C62828", "900": "B71C1C", "A100": "FF8A80", "A200": "FF5252",
        "A400": "FF1744", "A700": "D50000"
    }

    PINK = {
        "50": "FCE4EC", "100": "F8BBD0", "200": "F48FB1", "300": "F06292",
        "400": "EC407A", "500": "E91E63", "600": "D81B60", "700": "C2185B",
        "800": "AD1457", "900": "880E4F", "A100": "FF80AB", "A200": "FF4081",
        "A400": "F50057", "A700": "C51162"
    }

    PURPLE = {
        "50": "F3E5F5", "100": "E1BEE7", "200": "CE93D8", "300": "BA68C8",
        "400": "AB47BC", "500": "9C27B0", "600": "8E24AA", "700": "7B1FA2",
        "800": "6A1B9A", "900": "4A148C", "A100": "EA80FC", "A200": "E040FB",
        "A400": "D500F9", "A700": "AA00FF"
    }

    DEEP_PURPLE = {
        "50": "EDE7F6", "100": "D1C4E9", "200": "B39DDB", "300": "9575CD",
        "400": "7E57C2", "500": "673AB7", "600": "5E35B1", "700": "512DA8",
        "800": "4527A0", "900": "311B92", "A100": "B388FF", "A200": "7C4DFF",
        "A400": "651FFF", "A700": "6200EA"
    }

    INDIGO = {
        "50": "E8EAF6", "100": "C5CAE9", "200": "9FA8DA", "300": "7986CB",
        "400": "5C6BC0", "500": "3F51B5", "600": "3949AB", "700": "303F9F",
        "800": "283593", "900": "1A237E", "A100": "8C9EFF", "A200": "536DFE",
        "A400": "3D5AFE", "A700": "304FFE"
    }

    BLUE = {
        "50": "E3F2FD", "100": "BBDEFB", "200": "90CAF9", "300": "64B5F6",
        "400": "42A5F5", "500": "2196F3", "600": "1E88E5", "700": "1976D2",
        "800": "1565C0", "900": "0D47A1", "A100": "82B1FF", "A200": "448AFF",
        "A400": "2979FF", "A700": "2962FF"
    }

    LIGHT_BLUE = {
        "50": "E1F5FE", "100": "B3E5FC", "200": "81D4FA", "300": "4FC3F7",
        "400": "29B6F6", "500": "03A9F4", "600": "039BE5", "700": "0288D1",
        "800": "0277BD", "900": "01579B", "A100": "80D8FF", "A200": "40C4FF",
        "A400": "00B0FF", "A700": "0091EA"
    }

    CYAN = {
        "50": "E0F7FA", "100": "B2EBF2", "200": "80DEEA", "300": "4DD0E1",
        "400": "26C6DA", "500": "00BCD4", "600": "00ACC1", "700": "0097A7",
        "800": "00838F", "900": "006064", "A100": "84FFFF", "A200": "18FFFF",
        "A400": "00E5FF", "A700": "00B8D4"
    }

    TEAL = {
        "50": "E0F2F1", "100": "B2DFDB", "200": "80CBC4", "300": "4DB6AC",
        "400": "26A69A", "500": "009688", "600": "00897B", "700": "00796B",
        "800": "00695C", "900": "004D40", "A100": "A7FFEB", "A200": "64FFDA",
        "A400": "1DE9B6", "A700": "00BFA5"
    }

    GREEN = {
        "50": "E8F5E9", "100": "C8E6C9", "200": "A5D6A7", "300": "81C784",
        "400": "66BB6A", "500": "4CAF50", "600": "43A047", "700": "388E3C",
        "800": "2E7D32", "900": "1B5E20", "A100": "B9F6CA", "A200": "69F0AE",
        "A400": "00E676", "A700": "00C853"
    }

    LIGHT_GREEN = {
        "50": "F1F8E9", "100": "DCEDC8", "200": "C5E1A5", "300": "AED581",
        "400": "9CCC65", "500": "8BC34A", "600": "7CB342", "700": "689F38",
        "800": "558B2F", "900": "33691E", "A100": "CCFF90", "A200": "B2FF59",
        "A400": "76FF03", "A700": "64DD17"
    }

    LIME = {
        "50": "F9FBE7", "100": "F0F4C3", "200": "E6EE9C", "300": "DCE775",
        "400": "D4E157", "500": "CDDC39", "600": "C0CA33", "700": "AFB42B",
        "800": "9E9D24", "900": "827717", "A100": "F4FF81", "A200": "EEFF41",
        "A400": "C6FF00", "A700": "AEEA00"
    }

    AMBER = {
        "50": "FFF8E1", "100": "FFECB3", "200": "FFE082", "300": "FFD54F",
        "400": "FFCA28", "500": "FFC107", "600": "FFB300", "700": "FFA000",
        "800": "FF8F00", "900": "FF6F00", "A100": "FFE57F", "A200": "FFD740",
        "A400": "FFC400", "A700": "FFAB00"
    }

    ORANGE = {
        "50": "FFF3E0", "100": "FFE0B2", "200": "FFCC80", "300": "FFB74D",
        "400": "FFA726", "500": "FF9800", "600": "FB8C00", "700": "F57C00",
        "800": "EF6C00", "900": "E65100", "A100": "FFD180", "A200": "FFAB40",
        "A400": "FF9100", "A700": "FF6D00"
    }

    DEEP_ORANGE = {
        "50": "FBE9E7", "100": "FFCCBC", "200": "FFAB91", "300": "FF8A65",
        "400": "FF7043", "500": "FF5722", "600": "F4511E", "700": "E64A19",
        "800": "D84315", "900": "BF360C", "A100": "FF9E80", "A200": "FF6E40",
        "A400": "FF3D00", "A700": "DD2C00"
    }

    BROWN = {
        "50": "EFEBE9", "100": "D7CCC8", "200": "BCAAA4", "300": "A1887F",
        "400": "8D6E63", "500": "795548", "600": "6D4C41", "700": "5D4037",
        "800": "4E342E", "900": "3E2723"
    }

    GREY = {
        "50": "FAFAFA", "100": "F5F5F5", "200": "EEEEEE", "300": "E0E0E0",
        "400": "BDBDBD", "500": "9E9E9E", "600": "757575", "700": "616161",
        "800": "424242", "900": "212121"
    }

    BLUE_GREY = {
        "50": "ECEFF1", "100": "CFD8DC", "200": "B0BEC5", "300": "90A4AE",
        "400": "78909C", "500": "607D8B", "600": "546E7A", "700": "455A64",
        "800": "37474F", "900": "263238"
    }

@dataclass
class MaterialTheme:
    """Material Design Theme Configuration"""
    name: str
    primary_color: str
    primary_variant: str
    secondary_color: str
    secondary_variant: str
    background: str
    surface: str
    error: str
    on_primary: str
    on_secondary: str
    on_background: str
    on_surface: str
    on_error: str
    typography: Dict[str, Dict]
    elevation_shadows: List[Dict]
    spacing: Dict[str, float]
    corner_radius: float

class MaterialDesignThemes:
    """Pre-configured Material Design Themes"""

    @staticmethod
    def get_material_you_theme(seed_color: str) -> MaterialTheme:
        """Generate Material You dynamic theme from seed color"""
        # Convert hex to RGB
        rgb = tuple(int(seed_color[i:i+2], 16) for i in (0, 2, 4))
        h, l, s = colorsys.rgb_to_hls(rgb[0]/255, rgb[1]/255, rgb[2]/255)

        # Generate tonal palette
        primary = seed_color
        primary_variant = MaterialDesignThemes._adjust_color(seed_color, -0.2)
        secondary = MaterialDesignThemes._rotate_hue(seed_color, 30)
        secondary_variant = MaterialDesignThemes._adjust_color(secondary, -0.2)

        return MaterialTheme(
            name="Material You Dynamic",
            primary_color=primary,
            primary_variant=primary_variant,
            secondary_color=secondary,
            secondary_variant=secondary_variant,
            background="#FFFFFF",
            surface="#F5F5F5",
            error="#B00020",
            on_primary="#FFFFFF",
            on_secondary="#000000",
            on_background="#000000",
            on_surface="#000000",
            on_error="#FFFFFF",
            typography=MaterialDesignThemes._get_typography(),
            elevation_shadows=MaterialDesignThemes._get_shadows(),
            spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
            corner_radius=0.25
        )

    @staticmethod
    def _adjust_color(hex_color: str, factor: float) -> str:
        """Adjust color brightness"""
        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        h, l, s = colorsys.rgb_to_hls(rgb[0]/255, rgb[1]/255, rgb[2]/255)
        l = max(0, min(1, l + factor))
        r, g, b = colorsys.hls_to_rgb(h, l, s)
        return f"{int(r*255):02x}{int(g*255):02x}{int(b*255):02x}"

    @staticmethod
    def _rotate_hue(hex_color: str, degrees: float) -> str:
        """Rotate color hue"""
        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        h, l, s = colorsys.rgb_to_hls(rgb[0]/255, rgb[1]/255, rgb[2]/255)
        h = (h + degrees/360) % 1
        r, g, b = colorsys.hls_to_rgb(h, l, s)
        return f"{int(r*255):02x}{int(g*255):02x}{int(b*255):02x}"

    @staticmethod
    def _get_typography() -> Dict:
        """Get Material Design typography scale"""
        return {
            "h1": {"size": 96, "weight": "light", "spacing": -1.5},
            "h2": {"size": 60, "weight": "light", "spacing": -0.5},
            "h3": {"size": 48, "weight": "regular", "spacing": 0},
            "h4": {"size": 34, "weight": "regular", "spacing": 0.25},
            "h5": {"size": 24, "weight": "regular", "spacing": 0},
            "h6": {"size": 20, "weight": "medium", "spacing": 0.15},
            "subtitle1": {"size": 16, "weight": "regular", "spacing": 0.15},
            "subtitle2": {"size": 14, "weight": "medium", "spacing": 0.1},
            "body1": {"size": 16, "weight": "regular", "spacing": 0.5},
            "body2": {"size": 14, "weight": "regular", "spacing": 0.25},
            "button": {"size": 14, "weight": "medium", "spacing": 1.25},
            "caption": {"size": 12, "weight": "regular", "spacing": 0.4},
            "overline": {"size": 10, "weight": "regular", "spacing": 1.5}
        }

    @staticmethod
    def _get_shadows() -> List[Dict]:
        """Get Material Design elevation shadows"""
        return [
            {"elevation": 0, "shadow": None},
            {"elevation": 1, "shadow": "0px 2px 1px -1px rgba(0,0,0,0.2)"},
            {"elevation": 2, "shadow": "0px 3px 1px -2px rgba(0,0,0,0.2)"},
            {"elevation": 3, "shadow": "0px 3px 3px -2px rgba(0,0,0,0.2)"},
            {"elevation": 4, "shadow": "0px 2px 4px -1px rgba(0,0,0,0.2)"},
            {"elevation": 6, "shadow": "0px 3px 5px -1px rgba(0,0,0,0.2)"},
            {"elevation": 8, "shadow": "0px 5px 5px -3px rgba(0,0,0,0.2)"},
            {"elevation": 12, "shadow": "0px 7px 8px -4px rgba(0,0,0,0.2)"},
            {"elevation": 16, "shadow": "0px 8px 10px -5px rgba(0,0,0,0.2)"},
            {"elevation": 24, "shadow": "0px 11px 15px -7px rgba(0,0,0,0.2)"}
        ]

    @staticmethod
    def _get_all_themes():
        """Get all predefined themes - called after class is fully defined"""
        return {
            "material_baseline": MaterialTheme(
                name="Material Baseline",
                primary_color="6200EE",
                primary_variant="3700B3",
                secondary_color="03DAC6",
                secondary_variant="018786",
                background="FFFFFF",
                surface="FFFFFF",
                error="B00020",
                on_primary="FFFFFF",
                on_secondary="000000",
                on_background="000000",
                on_surface="000000",
                on_error="FFFFFF",
                typography=MaterialDesignThemes._get_typography(),
                elevation_shadows=MaterialDesignThemes._get_shadows(),
                spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
                corner_radius=0.25
            ),

        "material_dark": MaterialTheme(
            name="Material Dark",
            primary_color="BB86FC",
            primary_variant="3700B3",
            secondary_color="03DAC6",
            secondary_variant="03DAC6",
            background="121212",
            surface="121212",
            error="CF6679",
            on_primary="000000",
            on_secondary="000000",
            on_background="FFFFFF",
            on_surface="FFFFFF",
            on_error="000000",
            typography=MaterialDesignThemes._get_typography(),
            elevation_shadows=MaterialDesignThemes._get_shadows(),
            spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
            corner_radius=0.25
        ),

        "google_blue": MaterialTheme(
            name="Google Blue",
            primary_color="4285F4",
            primary_variant="1967D2",
            secondary_color="EA4335",
            secondary_variant="C5221F",
            background="FFFFFF",
            surface="F8F9FA",
            error="EA4335",
            on_primary="FFFFFF",
            on_secondary="FFFFFF",
            on_background="202124",
            on_surface="202124",
            on_error="FFFFFF",
            typography=MaterialDesignThemes._get_typography(),
            elevation_shadows=MaterialDesignThemes._get_shadows(),
            spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
            corner_radius=0.5
        ),

        "spotify_green": MaterialTheme(
            name="Spotify Green",
            primary_color="1DB954",
            primary_variant="1AA34A",
            secondary_color="191414",
            secondary_variant="000000",
            background="191414",
            surface="282828",
            error="E22134",
            on_primary="FFFFFF",
            on_secondary="FFFFFF",
            on_background="FFFFFF",
            on_surface="FFFFFF",
            on_error="FFFFFF",
            typography=MaterialDesignThemes._get_typography(),
            elevation_shadows=MaterialDesignThemes._get_shadows(),
            spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
            corner_radius=0.5
        ),

        "notion_minimal": MaterialTheme(
            name="Notion Minimal",
            primary_color="000000",
            primary_variant="2F3437",
            secondary_color="0084FF",
            secondary_variant="0073E6",
            background="FFFFFF",
            surface="F7F6F3",
            error="EB5757",
            on_primary="FFFFFF",
            on_secondary="FFFFFF",
            on_background="37352F",
            on_surface="37352F",
            on_error="FFFFFF",
            typography=MaterialDesignThemes._get_typography(),
            elevation_shadows=MaterialDesignThemes._get_shadows(),
            spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
            corner_radius=0.2
        )
        }
    
    # Cache for themes
    _THEMES_CACHE = None
    
    @classmethod
    def get_themes(cls):
        """Get all predefined themes (cached)"""
        if cls._THEMES_CACHE is None:
            cls._THEMES_CACHE = cls._get_all_themes()
        return cls._THEMES_CACHE

class MaterialDesignAdvisor:
    """Provides Material Design advice and best practices"""

    @staticmethod
    def get_color_advice(primary_color: str) -> Dict[str, Any]:
        """Get color combination advice based on primary color"""
        rgb = tuple(int(primary_color[i:i+2], 16) for i in (0, 2, 4))
        h, l, s = colorsys.rgb_to_hls(rgb[0]/255, rgb[1]/255, rgb[2]/255)

        advice = {
            "primary_color": primary_color,
            "color_psychology": MaterialDesignAdvisor._get_color_psychology(h),
            "recommended_combinations": MaterialDesignAdvisor._get_color_combinations(primary_color),
            "accessibility": MaterialDesignAdvisor._check_accessibility(primary_color),
            "usage_tips": MaterialDesignAdvisor._get_usage_tips(h, l, s)
        }

        return advice

    @staticmethod
    def _get_color_psychology(hue: float) -> str:
        """Get psychological impact of color"""
        if hue < 0.05 or hue > 0.95:  # Red
            return "Red conveys energy, passion, and urgency. Use for CTAs and important alerts."
        elif 0.05 <= hue < 0.15:  # Orange
            return "Orange suggests creativity and enthusiasm. Great for friendly, approachable brands."
        elif 0.15 <= hue < 0.2:  # Yellow
            return "Yellow evokes optimism and warmth. Use sparingly as accents to draw attention."
        elif 0.2 <= hue < 0.45:  # Green
            return "Green represents growth, harmony, and sustainability. Ideal for eco-friendly or financial themes."
        elif 0.45 <= hue < 0.55:  # Cyan
            return "Cyan/Teal suggests clarity and communication. Perfect for tech and healthcare."
        elif 0.55 <= hue < 0.75:  # Blue
            return "Blue conveys trust, stability, and professionalism. Most versatile for corporate use."
        elif 0.75 <= hue < 0.85:  # Purple
            return "Purple implies luxury, creativity, and wisdom. Excellent for premium or creative brands."
        else:  # Magenta
            return "Magenta/Pink expresses playfulness and innovation. Good for modern, bold designs."

    @staticmethod
    def _get_color_combinations(primary: str) -> List[Dict]:
        """Get recommended color combinations"""
        combinations = [
            {
                "name": "Complementary",
                "description": "High contrast, vibrant look",
                "colors": [primary, MaterialDesignThemes._rotate_hue(primary, 180)]
            },
            {
                "name": "Analogous",
                "description": "Harmonious and pleasing",
                "colors": [
                    MaterialDesignThemes._rotate_hue(primary, -30),
                    primary,
                    MaterialDesignThemes._rotate_hue(primary, 30)
                ]
            },
            {
                "name": "Triadic",
                "description": "Balanced and colorful",
                "colors": [
                    primary,
                    MaterialDesignThemes._rotate_hue(primary, 120),
                    MaterialDesignThemes._rotate_hue(primary, 240)
                ]
            },
            {
                "name": "Split Complementary",
                "description": "Vibrant with less tension",
                "colors": [
                    primary,
                    MaterialDesignThemes._rotate_hue(primary, 150),
                    MaterialDesignThemes._rotate_hue(primary, 210)
                ]
            }
        ]
        return combinations

    @staticmethod
    def _check_accessibility(color: str) -> Dict:
        """Check color accessibility"""
        # Calculate contrast ratios
        white_contrast = MaterialDesignAdvisor._get_contrast_ratio(color, "FFFFFF")
        black_contrast = MaterialDesignAdvisor._get_contrast_ratio(color, "000000")

        return {
            "wcag_aa_normal": white_contrast >= 4.5 or black_contrast >= 4.5,
            "wcag_aa_large": white_contrast >= 3 or black_contrast >= 3,
            "wcag_aaa_normal": white_contrast >= 7 or black_contrast >= 7,
            "wcag_aaa_large": white_contrast >= 4.5 or black_contrast >= 4.5,
            "best_text_color": "FFFFFF" if black_contrast > white_contrast else "000000",
            "contrast_white": round(white_contrast, 2),
            "contrast_black": round(black_contrast, 2)
        }

    @staticmethod
    def _get_contrast_ratio(color1: str, color2: str) -> float:
        """Calculate contrast ratio between two colors"""
        def get_luminance(hex_color):
            rgb = tuple(int(hex_color[i:i+2], 16)/255 for i in (0, 2, 4))
            rgb = [((c/12.92) if c <= 0.03928 else ((c+0.055)/1.055)**2.4) for c in rgb]
            return 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]

        l1 = get_luminance(color1)
        l2 = get_luminance(color2)

        return (max(l1, l2) + 0.05) / (min(l1, l2) + 0.05)

    @staticmethod
    def _get_usage_tips(h: float, l: float, s: float) -> List[str]:
        """Get usage tips based on color properties"""
        tips = []

        if l < 0.3:
            tips.append("Dark color - use for text or backgrounds with light content")
        elif l > 0.7:
            tips.append("Light color - best as background or with dark text overlay")
        else:
            tips.append("Medium brightness - versatile for both backgrounds and accents")

        if s > 0.7:
            tips.append("High saturation - use sparingly as accent color")
        elif s < 0.3:
            tips.append("Low saturation - suitable for large areas and backgrounds")

        tips.append("Follow 60-30-10 rule: 60% dominant, 30% secondary, 10% accent")
        tips.append("Maintain consistent color temperature throughout presentation")
        tips.append("Use color to create visual hierarchy and guide attention")

        return tips

    @staticmethod
    def get_layout_advice(slide_type: str) -> Dict[str, Any]:
        """Get Material Design layout advice for specific slide types"""
        layouts = {
            "title": {
                "grid": "Center aligned with generous padding",
                "spacing": "Use 8dp grid system",
                "typography": "Display large (h1/h2) for title, h4/h5 for subtitle",
                "tips": [
                    "Keep title concise - max 2 lines",
                    "Subtitle should support, not repeat the title",
                    "Add breathing room - don't fear white space"
                ]
            },
            "content": {
                "grid": "12 column grid with responsive breakpoints",
                "spacing": "16dp between elements, 24dp margins",
                "typography": "h4 for headers, body1 for content",
                "tips": [
                    "Limit to 3-5 bullet points per slide",
                    "Use progressive disclosure for complex info",
                    "Align elements to create clear scan lines"
                ]
            },
            "image": {
                "grid": "Full bleed or contained with consistent padding",
                "spacing": "Minimum 8dp padding from edges",
                "typography": "Caption in subtitle2 or caption style",
                "tips": [
                    "Use high-quality images (min 1920x1080)",
                    "Apply consistent image treatment (filters/overlays)",
                    "Ensure text contrast over images"
                ]
            },
            "data": {
                "grid": "Structured grid for charts and tables",
                "spacing": "Clear separation between data sets",
                "typography": "Clear hierarchy with h5 for titles, body2 for data",
                "tips": [
                    "Simplify data - highlight key insights",
                    "Use color purposefully to encode meaning",
                    "Add context with annotations"
                ]
            }
        }

        return layouts.get(slide_type, layouts["content"])

# Extension to the main server class
class MaterialDesignPowerPointExtension:
    """Extension for Material Design themed presentations"""

    def __init__(self, server):
        self.server = server
        self.themes = MaterialDesignThemes()
        self.advisor = MaterialDesignAdvisor()

    def add_material_tools(self):
        """Add Material Design specific tools to the server"""

        @self.server.server.list_tools()
        async def handle_list_material_tools() -> list[types.Tool]:
            """List Material Design tools"""
            return [
                types.Tool(
                    name="apply_material_theme",
                    description="Apply a Material Design theme to presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "theme_name": {
                                "type": "string",
                                "description": "Theme: material_baseline, material_dark, google_blue, spotify_green, notion_minimal, or custom"
                            },
                            "seed_color": {
                                "type": "string",
                                "description": "Hex color for Material You dynamic theme"
                            },
                            "dark_mode": {"type": "boolean", "default": False},
                            "high_contrast": {"type": "boolean", "default": False}
                        },
                        "required": ["presentation_id", "theme_name"]
                    }
                ),

                types.Tool(
                    name="get_material_color_palette",
                    description="Get Material Design color palette suggestions",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "base_color": {"type": "string", "description": "Hex color code"},
                            "palette_type": {
                                "type": "string",
                                "description": "Palette type: complementary, analogous, triadic, monochromatic"
                            }
                        },
                        "required": ["base_color"]
                    }
                ),

                types.Tool(
                    name="get_design_advice",
                    description="Get Material Design advice for presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "advice_type": {
                                "type": "string",
                                "description": "Type: color, typography, layout, spacing, animation, accessibility"
                            },
                            "context": {
                                "type": "string",
                                "description": "Context: corporate, educational, creative, technical, marketing"
                            }
                        },
                        "required": ["advice_type"]
                    }
                ),

                types.Tool(
                    name="apply_material_layout",
                    description="Apply Material Design layout patterns",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "layout_pattern": {
                                "type": "string",
                                "description": "Pattern: hero, cards, list, dashboard, feature, testimonial"
                            },
                            "responsive": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id", "slide_index", "layout_pattern"]
                    }
                ),

                types.Tool(
                    name="add_material_components",
                    description="Add Material Design components to slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "component_type": {
                                "type": "string",
                                "description": "Component: fab, chip, card, badge, snackbar, stepper"
                            },
                            "properties": {"type": "object"}
                        },
                        "required": ["presentation_id", "slide_index", "component_type"]
                    }
                ),

                types.Tool(
                    name="check_accessibility",
                    description="Check presentation accessibility against WCAG standards",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "standard": {
                                "type": "string",
                                "description": "Standard: WCAG_AA, WCAG_AAA",
                                "default": "WCAG_AA"
                            }
                        },
                        "required": ["presentation_id"]
                    }
                ),

                types.Tool(
                    name="optimize_for_device",
                    description="Optimize presentation for specific devices",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "device_type": {
                                "type": "string",
                                "description": "Device: desktop, tablet, mobile, tv"
                            },
                            "orientation": {
                                "type": "string",
                                "description": "Orientation: landscape, portrait",
                                "default": "landscape"
                            }
                        },
                        "required": ["presentation_id", "device_type"]
                    }
                ),

                types.Tool(
                    name="generate_style_guide",
                    description="Generate a style guide slide for the presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "include_colors": {"type": "boolean", "default": True},
                            "include_typography": {"type": "boolean", "default": True},
                            "include_spacing": {"type": "boolean", "default": True},
                            "include_components": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id"]
                    }
                ),

                types.Tool(
                    name="apply_material_animations",
                    description="Apply Material Design motion principles",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "animation_style": {
                                "type": "string",
                                "description": "Style: standard, emphasized, expressive"
                            },
                            "easing": {
                                "type": "string",
                                "description": "Easing: standard, decelerated, accelerated",
                                "default": "standard"
                            },
                            "duration_scale": {"type": "number", "default": 1.0}
                        },
                        "required": ["presentation_id", "animation_style"]
                    }
                ),

                types.Tool(
                    name="create_mood_board",
                    description="Create a mood board slide with Material Design elements",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "theme": {
                                "type": "string",
                                "description": "Theme mood: professional, playful, elegant, minimal, bold"
                            },
                            "industry": {
                                "type": "string",
                                "description": "Industry context for appropriate styling"
                            }
                        },
                        "required": ["presentation_id", "theme"]
                    }
                )
            ]

    async def apply_material_theme(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Apply Material Design theme to presentation"""
        pres_id = args["presentation_id"]
        if pres_id not in self.server.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.server.presentations[pres_id]
        theme_name = args["theme_name"]

        # Get or create theme
        if theme_name == "custom" and "seed_color" in args:
            theme = MaterialDesignThemes.get_material_you_theme(args["seed_color"])
        else:
            theme = MaterialDesignThemes.get_themes().get(theme_name)
            if not theme:
                raise ValueError(f"Unknown theme: {theme_name}")

        # Apply dark mode if requested
        if args.get("dark_mode", False):
            theme = self._apply_dark_mode(theme)

        # Apply high contrast if requested
        if args.get("high_contrast", False):
            theme = self._apply_high_contrast(theme)

        # Apply theme to all slides
        for slide in prs.slides:
            self._apply_theme_to_slide(slide, theme)

        # Store theme metadata
        if not hasattr(prs, 'material_theme'):
            prs.material_theme = theme

        return [types.TextContent(
            type="text",
            text=f"Applied Material Design theme '{theme.name}' to presentation"
        )]

    def _apply_theme_to_slide(self, slide, theme: MaterialTheme):
        """Apply Material theme to a single slide"""
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(theme.background))

        # Apply to text elements
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(*self._hex_to_rgb(theme.on_background))

                        # Apply typography based on level
                        if shape == slide.shapes.title:
                            run.font.size = Pt(theme.typography["h3"]["size"])
                        else:
                            run.font.size = Pt(theme.typography["body1"]["size"])

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _apply_dark_mode(self, theme: MaterialTheme) -> MaterialTheme:
        """Convert theme to dark mode"""
        dark_theme = MaterialTheme(
            name=f"{theme.name} Dark",
            primary_color=self._lighten_color(theme.primary_color, 0.2),
            primary_variant=theme.primary_variant,
            secondary_color=self._lighten_color(theme.secondary_color, 0.2),
            secondary_variant=theme.secondary_variant,
            background="121212",
            surface="1E1E1E",
            error="CF6679",
            on_primary="000000",
            on_secondary="000000",
            on_background="FFFFFF",
            on_surface="FFFFFF",
            on_error="000000",
            typography=theme.typography,
            elevation_shadows=theme.elevation_shadows,
            spacing=theme.spacing,
            corner_radius=theme.corner_radius
        )
        return dark_theme

    def _apply_high_contrast(self, theme: MaterialTheme) -> MaterialTheme:
        """Apply high contrast to theme"""
        # Increase contrast by making lights lighter and darks darker
        theme.on_background = "000000" if theme.background > "7F7F7F" else "FFFFFF"
        theme.on_surface = "000000" if theme.surface > "7F7F7F" else "FFFFFF"
        return theme

    def _lighten_color(self, hex_color: str, factor: float) -> str:
        """Lighten a color"""
        return MaterialDesignThemes._adjust_color(hex_color, factor)

    async def get_design_advice(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Get Material Design advice"""
        advice_type = args["advice_type"]
        context = args.get("context", "general")

        advice = {
            "color": self._get_color_advice(context),
            "typography": self._get_typography_advice(context),
            "layout": self._get_layout_advice(context),
            "spacing": self._get_spacing_advice(context),
            "animation": self._get_animation_advice(context),
            "accessibility": self._get_accessibility_advice(context)
        }

        result = advice.get(advice_type, "Unknown advice type")

        return [types.TextContent(
            type="text",
            text=json.dumps(result, indent=2)
        )]

    def _get_color_advice(self, context: str) -> Dict:
        """Get color advice for context"""
        return {
            "principles": [
                "Use color purposefully to communicate meaning",
                "Maintain sufficient contrast (WCAG AA minimum)",
                "Limit palette to 3-4 colors maximum",
                "Use semantic colors consistently (error, success, warning)"
            ],
            "context_specific": {
                "corporate": "Use professional blues and grays, minimal accent colors",
                "educational": "Bright, engaging colors with good contrast for readability",
                "creative": "Bold, expressive colors with unexpected combinations",
                "technical": "Dark themes with syntax highlighting colors",
                "marketing": "Brand colors with emotional appeal"
            }.get(context, "Choose colors that align with your message"),
            "tips": [
                "Test colors with color blindness simulators",
                "Use opacity for hierarchy without adding colors",
                "Consider cultural color associations for global audiences"
            ]
        }

    def _get_typography_advice(self, context: str) -> Dict:
        """Get typography advice"""
        return {
            "hierarchy": {
                "title": "Use display or h1-h3 for impact",
                "subtitle": "Use h4-h6 for section headers",
                "body": "Use body1 (16pt) for readability",
                "caption": "Use caption or overline for metadata"
            },
            "guidelines": [
                "Maintain consistent type scale throughout",
                "Limit to 2 font families maximum",
                "Use weight and size for hierarchy, not just color",
                "Ensure 60-75 characters per line for body text"
            ],
            "context_specific": {
                "corporate": "Professional serif or clean sans-serif",
                "educational": "Highly legible sans-serif with clear distinction",
                "creative": "Expressive display fonts with personality",
                "technical": "Monospace for code, clean sans for explanations",
                "marketing": "Brand fonts with emotional resonance"
            }.get(context, "Choose fonts that enhance readability")
        }

    def _get_layout_advice(self, context: str) -> Dict:
        """Get layout advice"""
        return {
            "grid": "Use 8dp or 12-column grid for consistency",
            "alignment": "Align elements to create clear relationships",
            "whitespace": "Use generous spacing to improve focus",
            "visual_hierarchy": [
                "Size - Larger elements draw attention first",
                "Color - Bright or contrasting colors stand out",
                "Position - Top-left gets attention in LTR languages",
                "Space - Isolated elements get more attention"
            ],
            "patterns": {
                "F-pattern": "For text-heavy content",
                "Z-pattern": "For minimal content with CTA",
                "Card-based": "For organizing related content groups",
                "Hero": "For impactful opening statements"
            }
        }

    def _get_spacing_advice(self, context: str) -> Dict:
        """Get spacing advice"""
        return {
            "system": "Use 8dp grid system (8, 16, 24, 32, 40, 48, 56, 64)",
            "margins": {
                "mobile": "16dp",
                "tablet": "24dp",
                "desktop": "24-48dp"
            },
            "padding": {
                "dense": "4-8dp for compact layouts",
                "comfortable": "16dp standard padding",
                "relaxed": "24-32dp for spacious feel"
            },
            "line_height": {
                "tight": "1.2 for headlines",
                "normal": "1.5 for body text",
                "loose": "1.75-2 for improved readability"
            }
        }

    def _get_animation_advice(self, context: str) -> Dict:
        """Get animation advice"""
        return {
            "principles": [
                "Make it responsive - animations respond to user input",
                "Natural - follow physics and feel intuitive",
                "Aware - spatial and hierarchical relationships",
                "Intentional - guide focus and provide feedback"
            ],
            "duration": {
                "small": "100-200ms for small transitions",
                "medium": "200-300ms for medium complexity",
                "large": "300-400ms for complex or large animations",
                "extra_large": "400-500ms for dramatic effects"
            },
            "easing": {
                "standard": "cubic-bezier(0.4, 0.0, 0.2, 1)",
                "decelerate": "cubic-bezier(0.0, 0.0, 0.2, 1)",
                "accelerate": "cubic-bezier(0.4, 0.0, 1, 1)"
            }
        }

    def _get_accessibility_advice(self, context: str) -> Dict:
        """Get accessibility advice"""
        return {
            "contrast": {
                "WCAG_AA": "4.5:1 for normal text, 3:1 for large text",
                "WCAG_AAA": "7:1 for normal text, 4.5:1 for large text"
            },
            "guidelines": [
                "Provide text alternatives for images",
                "Don't rely solely on color to convey information",
                "Ensure keyboard navigation is possible",
                "Use semantic markup for screen readers",
                "Provide captions for videos",
                "Test with accessibility tools"
            ],
            "inclusive_design": [
                "Consider color blindness (8% of men, 0.5% of women)",
                "Account for low vision users",
                "Design for one-handed use on mobile",
                "Provide multiple ways to access information"
            ]
        }

    async def generate_style_guide(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Generate a comprehensive style guide slide"""
        pres_id = args["presentation_id"]
        if pres_id not in self.server.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.server.presentations[pres_id]

        # Create style guide slide
        slide_layout = prs.slide_layouts[5]  # Title only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Style Guide"

        # Add color palette
        if args.get("include_colors", True):
            self._add_color_palette(slide)

        # Add typography samples
        if args.get("include_typography", True):
            self._add_typography_samples(slide)

        # Add spacing guide
        if args.get("include_spacing", True):
            self._add_spacing_guide(slide)

        # Add component samples
        if args.get("include_components", True):
            self._add_component_samples(slide)

        return [types.TextContent(
            type="text",
            text="Generated comprehensive style guide slide"
        )]

    def _add_color_palette(self, slide):
        """Add color palette to style guide"""
        left = Inches(1)
        top = Inches(1.5)
        size = Inches(0.75)

        colors = [
            ("Primary", "2196F3"),
            ("Secondary", "FF5722"),
            ("Success", "4CAF50"),
            ("Warning", "FFC107"),
            ("Error", "F44336"),
            ("Background", "FAFAFA")
        ]

        for i, (name, color) in enumerate(colors):
            # Color swatch
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + (i % 3) * Inches(3),
                top + (i // 3) * Inches(1.5),
                size,
                size
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(color))

            # Label
            text_box = slide.shapes.add_textbox(
                left + (i % 3) * Inches(3),
                top + (i // 3) * Inches(1.5) + size + Inches(0.1),
                size,
                Inches(0.3)
            )
            text_box.text = f"{name}\\n#{color}"
            text_box.text_frame.paragraphs[0].font.size = Pt(10)

    def _add_typography_samples(self, slide):
        """Add typography samples to style guide"""
        left = Inches(1)
        top = Inches(3.5)

        samples = [
            ("Heading 1", 32, False),
            ("Heading 2", 24, False),
            ("Body", 14, False),
            ("Caption", 12, True)
        ]

        for i, (label, size, italic) in enumerate(samples):
            text_box = slide.shapes.add_textbox(
                left,
                top + i * Inches(0.5),
                Inches(8),
                Inches(0.4)
            )
            p = text_box.text_frame.paragraphs[0]
            p.text = f"{label}: The quick brown fox jumps over the lazy dog"
            p.font.size = Pt(size)
            p.font.italic = italic

    def _add_spacing_guide(self, slide):
        """Add spacing guide to style guide"""
        # Implementation for spacing visualization
        pass

    def _add_component_samples(self, slide):
        """Add component samples to style guide"""
        # Implementation for component samples
        pass

# Export classes for unified server
__all__ = [
    'MaterialColorPalette',
    'MaterialTheme',
    'MaterialDesignThemes',
    'MaterialDesignAdvisor',
    'MaterialDesignPowerPointExtension'
]
