#!/usr/bin/env python3
"""
Markdown to PowerPoint Automation System with Visual Tags
Supports custom markdown syntax for advanced PowerPoint features
"""

import re
import json
import yaml
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass, field
from enum import Enum
import asyncio
from pathlib import Path
import hashlib
import tempfile
import urllib.request
import base64
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
try:
    import markdown
except ImportError:
    print("Warning: markdown module not found. Install with: pip install markdown")
from PIL import Image

class SlideType(Enum):
    """Supported slide types"""
    TITLE = "title"
    SECTION = "section"
    CONTENT = "content"
    TWO_COLUMN = "two_column"
    IMAGE = "image"
    CHART = "chart"
    TABLE = "table"
    QUOTE = "quote"
    CODE = "code"
    TIMELINE = "timeline"
    COMPARISON = "comparison"
    TEAM = "team"
    BLANK = "blank"

@dataclass
class SlideConfig:
    """Configuration for a single slide"""
    type: SlideType
    title: Optional[str] = None
    subtitle: Optional[str] = None
    content: List[Any] = field(default_factory=list)
    layout: Optional[str] = None
    background: Optional[str] = None
    transition: Optional[str] = None
    notes: Optional[str] = None
    animation: Optional[str] = None
    theme_override: Optional[Dict] = None
    metadata: Dict = field(default_factory=dict)

@dataclass
class PresentationConfig:
    """Global presentation configuration"""
    title: str
    author: Optional[str] = None
    theme: str = "default"
    aspect_ratio: str = "16:9"
    footer: Optional[str] = None
    slide_numbers: bool = True
    date: bool = False
    company: Optional[str] = None
    logo: Optional[str] = None
    color_scheme: Optional[Dict] = None
    font_family: Optional[str] = None
    metadata: Dict = field(default_factory=dict)

class MarkdownPresentationParser:
    """Parse markdown with custom PowerPoint tags"""

    # Regex patterns for custom tags
    PATTERNS = {
        # Slide markers
        'new_slide': r'^---\\s*$',
        'slide_type': r'^<!--\\s*slide:\\s*(\\w+)\\s*-->',

        # Visual features
        'background': r'<!--\\s*background:\\s*([^-]+?)\\s*-->',
        'background_image': r'<!--\\s*bg-image:\\s*([^-]+?)\\s*-->',
        'background_video': r'<!--\\s*bg-video:\\s*([^-]+?)\\s*-->',
        'transition': r'<!--\\s*transition:\\s*(\\w+)(?:\\s+(\\d+))?\\s*-->',
        'animation': r'<!--\\s*animate:\\s*(\\w+)(?:\\s+(\\d+))?\\s*-->',
        'layout': r'<!--\\s*layout:\\s*(\\w+)\\s*-->',
        'theme': r'<!--\\s*theme:\\s*(\\w+)\\s*-->',
        'notes': r'<!--\\s*notes:\\s*(.*?)\\s*-->',

        # Content elements
        'image': r'!\\[([^\\]]*)\\]\\(([^)]+)\\)(?:\\{([^}]+)\\})?',
        'chart': r'```chart\\s*\\n(.*?)```',
        'table': r'```table\\s*\\n(.*?)```',
        'code': r'```(\\w+)\\s*\\n(.*?)```',
        'math': r'\\$\\$(.*?)\\$\\$',
        'timeline': r'```timeline\\s*\\n(.*?)```',
        'mermaid': r'```mermaid\\s*\\n(.*?)```',

        # Advanced features
        'columns': r':::\\s*columns\\s*\\n(.*?):::',
        'speaker': r'<!--\\s*speaker:\\s*(.*?)\\s*-->',
        'timer': r'<!--\\s*timer:\\s*(\\d+)\\s*-->',
        'poll': r'<!--\\s*poll:\\s*(.*?)\\s*-->',
        'qr': r'<!--\\s*qr:\\s*(.*?)\\s*-->',

        # Styling
        'color': r'\\{color:([^}]+)\\}',
        'size': r'\\{size:([^}]+)\\}',
        'align': r'\\{align:([^}]+)\\}',
        'font': r'\\{font:([^}]+)\\}',
        'box': r':::\\s*box(?:\\s+(\\w+))?\\s*\\n(.*?):::',

        # Variables and includes
        'variable': r'\\{\\{(\\w+)\\}\\}',
        'include': r'<!--\\s*include:\\s*([^-]+?)\\s*-->',

        # Metadata
        'meta': r'^---\\s*\\n(.*?)\\n---\\s*$',
    }

    def __init__(self):
        self.slides: List[SlideConfig] = []
        self.presentation_config: Optional[PresentationConfig] = None
        self.variables: Dict[str, Any] = {}
        self.current_slide: Optional[SlideConfig] = None

    def parse(self, markdown_content: str) -> Tuple[PresentationConfig, List[SlideConfig]]:
        """Parse markdown content into presentation configuration"""

        # Extract metadata if present
        self._parse_metadata(markdown_content)

        # Split into slides
        slides_raw = re.split(self.PATTERNS['new_slide'], markdown_content, flags=re.MULTILINE)

        for slide_raw in slides_raw:
            if slide_raw.strip():
                self._parse_slide(slide_raw)

        return self.presentation_config, self.slides

    def _parse_metadata(self, content: str):
        """Parse YAML front matter"""
        meta_match = re.search(self.PATTERNS['meta'], content, re.MULTILINE | re.DOTALL)
        if meta_match:
            try:
                metadata = yaml.safe_load(meta_match.group(1))
                self.presentation_config = PresentationConfig(
                    title=metadata.get('title', 'Presentation'),
                    author=metadata.get('author'),
                    theme=metadata.get('theme', 'default'),
                    aspect_ratio=metadata.get('aspect_ratio', '16:9'),
                    footer=metadata.get('footer'),
                    slide_numbers=metadata.get('slide_numbers', True),
                    date=metadata.get('date', False),
                    company=metadata.get('company'),
                    logo=metadata.get('logo'),
                    color_scheme=metadata.get('colors'),
                    font_family=metadata.get('font'),
                    metadata=metadata
                )
                self.variables = metadata.get('variables', {})
            except yaml.YAMLError:
                pass

        if not self.presentation_config:
            self.presentation_config = PresentationConfig(title="Presentation")

    def _parse_slide(self, slide_content: str):
        """Parse individual slide content"""
        slide = SlideConfig(type=SlideType.CONTENT)

        # Detect slide type
        type_match = re.search(self.PATTERNS['slide_type'], slide_content)
        if type_match:
            try:
                slide.type = SlideType(type_match.group(1))
            except ValueError:
                slide.type = SlideType.CONTENT
        else:
            # Auto-detect slide type based on content
            slide.type = self._detect_slide_type(slide_content)

        # Parse slide properties
        self._parse_slide_properties(slide, slide_content)

        # Parse content based on slide type
        self._parse_slide_content(slide, slide_content)

        self.slides.append(slide)

    def _detect_slide_type(self, content: str) -> SlideType:
        """Auto-detect slide type from content patterns"""
        lines = content.strip().split('\\n')

        # Title slide: single # heading with optional subtitle
        if len(lines) <= 3 and lines[0].startswith('# '):
            return SlideType.TITLE

        # Section slide: ## heading
        if lines[0].startswith('## ') and len(lines) <= 2:
            return SlideType.SECTION

        # Chart slide
        if '```chart' in content:
            return SlideType.CHART

        # Table slide
        if '```table' in content or '|' in content:
            return SlideType.TABLE

        # Code slide
        if re.search(r'```\\w+', content):
            return SlideType.CODE

        # Quote slide
        if content.strip().startswith('>'):
            return SlideType.QUOTE

        # Timeline slide
        if '```timeline' in content:
            return SlideType.TIMELINE

        # Two column if columns marker present
        if ':::columns' in content:
            return SlideType.TWO_COLUMN

        # Default to content slide
        return SlideType.CONTENT

    def _parse_slide_properties(self, slide: SlideConfig, content: str):
        """Parse slide properties from comments"""

        # Background
        bg_match = re.search(self.PATTERNS['background'], content)
        if bg_match:
            slide.background = bg_match.group(1).strip()

        bg_img_match = re.search(self.PATTERNS['background_image'], content)
        if bg_img_match:
            slide.metadata['bg_image'] = bg_img_match.group(1).strip()

        # Transition
        trans_match = re.search(self.PATTERNS['transition'], content)
        if trans_match:
            slide.transition = trans_match.group(1)
            if trans_match.group(2):
                slide.metadata['transition_duration'] = int(trans_match.group(2))

        # Animation
        anim_match = re.search(self.PATTERNS['animation'], content)
        if anim_match:
            slide.animation = anim_match.group(1)
            if anim_match.group(2):
                slide.metadata['animation_delay'] = int(anim_match.group(2))

        # Layout
        layout_match = re.search(self.PATTERNS['layout'], content)
        if layout_match:
            slide.layout = layout_match.group(1)

        # Speaker notes
        notes_match = re.search(self.PATTERNS['notes'], content)
        if notes_match:
            slide.notes = notes_match.group(1).strip()

        # Theme override
        theme_match = re.search(self.PATTERNS['theme'], content)
        if theme_match:
            slide.metadata['theme'] = theme_match.group(1)

    def _parse_slide_content(self, slide: SlideConfig, content: str):
        """Parse slide content based on type"""

        # Clean content from property comments
        clean_content = re.sub(r'<!--.*?-->', '', content, flags=re.DOTALL)

        if slide.type == SlideType.TITLE:
            self._parse_title_slide(slide, clean_content)
        elif slide.type == SlideType.SECTION:
            self._parse_section_slide(slide, clean_content)
        elif slide.type == SlideType.CONTENT:
            self._parse_content_slide(slide, clean_content)
        elif slide.type == SlideType.TWO_COLUMN:
            self._parse_two_column_slide(slide, clean_content)
        elif slide.type == SlideType.CHART:
            self._parse_chart_slide(slide, clean_content)
        elif slide.type == SlideType.TABLE:
            self._parse_table_slide(slide, clean_content)
        elif slide.type == SlideType.CODE:
            self._parse_code_slide(slide, clean_content)
        elif slide.type == SlideType.QUOTE:
            self._parse_quote_slide(slide, clean_content)
        elif slide.type == SlideType.TIMELINE:
            self._parse_timeline_slide(slide, clean_content)
        elif slide.type == SlideType.IMAGE:
            self._parse_image_slide(slide, clean_content)

    def _parse_title_slide(self, slide: SlideConfig, content: str):
        """Parse title slide content"""
        lines = content.strip().split('\\n')
        if lines:
            slide.title = lines[0].replace('# ', '').strip()
            if len(lines) > 1:
                slide.subtitle = '\\n'.join(lines[1:]).strip()

    def _parse_section_slide(self, slide: SlideConfig, content: str):
        """Parse section slide content"""
        lines = content.strip().split('\\n')
        if lines:
            slide.title = lines[0].replace('## ', '').strip()
            if len(lines) > 1:
                slide.subtitle = '\\n'.join(lines[1:]).strip()

    def _parse_content_slide(self, slide: SlideConfig, content: str):
        """Parse content slide with bullets"""
        lines = content.strip().split('\\n')

        # Extract title
        for i, line in enumerate(lines):
            if line.startswith('#'):
                slide.title = line.lstrip('#').strip()
                lines = lines[i+1:]
                break

        # Parse content
        slide.content = self._parse_markdown_list(lines)

    def _parse_markdown_list(self, lines: List[str]) -> List[Dict]:
        """Parse markdown list into structured format"""
        items = []
        current_item = None
        current_level = 0

        for line in lines:
            # Skip empty lines
            if not line.strip():
                continue

            # Detect list items
            match = re.match(r'^(\\s*)[-*+]\\s+(.+)', line)
            if match:
                indent = len(match.group(1))
                level = indent // 2
                text = match.group(2)

                items.append({
                    'text': text,
                    'level': level,
                    'type': 'bullet'
                })
            else:
                # Regular paragraph
                items.append({
                    'text': line.strip(),
                    'level': 0,
                    'type': 'paragraph'
                })

        return items

    def _parse_two_column_slide(self, slide: SlideConfig, content: str):
        """Parse two-column slide"""
        columns_match = re.search(self.PATTERNS['columns'], content, re.DOTALL)
        if columns_match:
            columns_content = columns_match.group(1)
            columns = columns_content.split('|||')

            slide.content = {
                'left': self._parse_column_content(columns[0]) if len(columns) > 0 else [],
                'right': self._parse_column_content(columns[1]) if len(columns) > 1 else []
            }

        # Extract title
        title_match = re.search(r'^#+\\s+(.+)$', content, re.MULTILINE)
        if title_match:
            slide.title = title_match.group(1)

    def _parse_column_content(self, content: str) -> List[Any]:
        """Parse content within a column"""
        return self._parse_markdown_list(content.strip().split('\\n'))

    def _parse_chart_slide(self, slide: SlideConfig, content: str):
        """Parse chart slide"""
        chart_match = re.search(self.PATTERNS['chart'], content, re.DOTALL)
        if chart_match:
            chart_data = yaml.safe_load(chart_match.group(1))
            slide.content = [{
                'type': 'chart',
                'chart_type': chart_data.get('type', 'column'),
                'data': chart_data.get('data', {}),
                'options': chart_data.get('options', {})
            }]

        # Extract title
        title_match = re.search(r'^#+\\s+(.+)$', content, re.MULTILINE)
        if title_match:
            slide.title = title_match.group(1)

    def _parse_table_slide(self, slide: SlideConfig, content: str):
        """Parse table slide"""
        table_match = re.search(self.PATTERNS['table'], content, re.DOTALL)

        if table_match:
            # Parse YAML table format
            table_data = yaml.safe_load(table_match.group(1))
            slide.content = [{
                'type': 'table',
                'headers': table_data.get('headers', []),
                'rows': table_data.get('rows', []),
                'style': table_data.get('style', 'default')
            }]
        else:
            # Try to parse markdown table
            table_data = self._parse_markdown_table(content)
            if table_data:
                slide.content = [table_data]

        # Extract title
        title_match = re.search(r'^#+\\s+(.+)$', content, re.MULTILINE)
        if title_match:
            slide.title = title_match.group(1)

    def _parse_markdown_table(self, content: str) -> Optional[Dict]:
        """Parse markdown table format"""
        lines = content.strip().split('\\n')
        table_lines = [line for line in lines if '|' in line]

        if len(table_lines) < 2:
            return None

        # Parse headers
        headers = [cell.strip() for cell in table_lines[0].split('|')[1:-1]]

        # Skip separator line
        rows = []
        for line in table_lines[2:]:
            row = [cell.strip() for cell in line.split('|')[1:-1]]
            rows.append(row)

        return {
            'type': 'table',
            'headers': headers,
            'rows': rows,
            'style': 'default'
        }

    def _parse_code_slide(self, slide: SlideConfig, content: str):
        """Parse code slide"""
        code_match = re.search(self.PATTERNS['code'], content, re.DOTALL)
        if code_match:
            slide.content = [{
                'type': 'code',
                'language': code_match.group(1),
                'code': code_match.group(2).strip()
            }]

        # Extract title
        title_match = re.search(r'^#+\\s+(.+)$', content, re.MULTILINE)
        if title_match:
            slide.title = title_match.group(1)

    def _parse_quote_slide(self, slide: SlideConfig, content: str):
        """Parse quote slide"""
        lines = content.strip().split('\\n')
        quote_lines = []
        author = None

        for line in lines:
            if line.startswith('>'):
                quote_lines.append(line[1:].strip())
            elif line.startswith('—') or line.startswith('-'):
                author = line[1:].strip()

        slide.content = [{
            'type': 'quote',
            'text': ' '.join(quote_lines),
            'author': author
        }]

    def _parse_timeline_slide(self, slide: SlideConfig, content: str):
        """Parse timeline slide"""
        timeline_match = re.search(self.PATTERNS['timeline'], content, re.DOTALL)
        if timeline_match:
            timeline_data = yaml.safe_load(timeline_match.group(1))
            slide.content = [{
                'type': 'timeline',
                'events': timeline_data.get('events', []),
                'style': timeline_data.get('style', 'horizontal')
            }]

        # Extract title
        title_match = re.search(r'^#+\\s+(.+)$', content, re.MULTILINE)
        if title_match:
            slide.title = title_match.group(1)

    def _parse_image_slide(self, slide: SlideConfig, content: str):
        """Parse image slide"""
        img_matches = re.finditer(self.PATTERNS['image'], content)
        images = []

        for match in img_matches:
            img_data = {
                'type': 'image',
                'alt': match.group(1),
                'src': match.group(2)
            }

            # Parse image attributes
            if match.group(3):
                attrs = match.group(3)
                attr_dict = {}
                for attr in attrs.split(','):
                    if '=' in attr:
                        key, value = attr.split('=', 1)
                        attr_dict[key.strip()] = value.strip()
                img_data['attributes'] = attr_dict

            images.append(img_data)

        slide.content = images

        # Extract title
        title_match = re.search(r'^#+\\s+(.+)$', content, re.MULTILINE)
        if title_match:
            slide.title = title_match.group(1)

class PowerPointGenerator:
    """Generate PowerPoint from parsed markdown"""

    def __init__(self):
        self.prs: Optional[Presentation] = None
        self.theme_colors: Dict[str, RGBColor] = {}
        self.default_font: str = "Calibri"

    def generate(self, config: PresentationConfig, slides: List[SlideConfig]) -> Presentation:
        """Generate PowerPoint presentation from parsed content"""

        # Create presentation
        self.prs = Presentation()

        # Apply global configuration
        self._apply_global_config(config)

        # Generate each slide
        for slide_config in slides:
            self._generate_slide(slide_config)

        return self.prs

    def _apply_global_config(self, config: PresentationConfig):
        """Apply global presentation settings"""

        # Set aspect ratio
        if config.aspect_ratio == "16:9":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(5.625)
        elif config.aspect_ratio == "4:3":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)

        # Set theme colors
        if config.color_scheme:
            self.theme_colors = {
                'primary': self._parse_color(config.color_scheme.get('primary', '000000')),
                'secondary': self._parse_color(config.color_scheme.get('secondary', '666666')),
                'accent': self._parse_color(config.color_scheme.get('accent', '0066CC')),
                'background': self._parse_color(config.color_scheme.get('background', 'FFFFFF')),
                'text': self._parse_color(config.color_scheme.get('text', '000000'))
            }

        # Set font
        if config.font_family:
            self.default_font = config.font_family

    def _parse_color(self, color_str: str) -> RGBColor:
        """Parse color string to RGBColor"""
        color_str = color_str.lstrip('#')
        return RGBColor(
            int(color_str[0:2], 16),
            int(color_str[2:4], 16),
            int(color_str[4:6], 16)
        )

    def _generate_slide(self, slide_config: SlideConfig):
        """Generate a single slide"""

        if slide_config.type == SlideType.TITLE:
            self._generate_title_slide(slide_config)
        elif slide_config.type == SlideType.SECTION:
            self._generate_section_slide(slide_config)
        elif slide_config.type == SlideType.CONTENT:
            self._generate_content_slide(slide_config)
        elif slide_config.type == SlideType.TWO_COLUMN:
            self._generate_two_column_slide(slide_config)
        elif slide_config.type == SlideType.CHART:
            self._generate_chart_slide(slide_config)
        elif slide_config.type == SlideType.TABLE:
            self._generate_table_slide(slide_config)
        elif slide_config.type == SlideType.CODE:
            self._generate_code_slide(slide_config)
        elif slide_config.type == SlideType.QUOTE:
            self._generate_quote_slide(slide_config)
        elif slide_config.type == SlideType.TIMELINE:
            self._generate_timeline_slide(slide_config)
        elif slide_config.type == SlideType.IMAGE:
            self._generate_image_slide(slide_config)
        else:
            self._generate_blank_slide(slide_config)

        # Apply slide-specific settings
        if self.prs.slides:
            current_slide = self.prs.slides[-1]
            self._apply_slide_settings(current_slide, slide_config)

    def _generate_title_slide(self, config: SlideConfig):
        """Generate title slide"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if config.subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = config.subtitle

    def _generate_section_slide(self, config: SlideConfig):
        """Generate section slide"""
        slide_layout = self.prs.slide_layouts[2] if len(self.prs.slide_layouts) > 2 else self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if config.subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = config.subtitle

    def _generate_content_slide(self, config: SlideConfig):
        """Generate content slide with bullets"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        # Add content
        if config.content and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            for item in config.content:
                if isinstance(item, dict):
                    p = text_frame.add_paragraph()
                    p.text = item.get('text', '')
                    p.level = item.get('level', 0)
                elif isinstance(item, str):
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0

    def _generate_two_column_slide(self, config: SlideConfig):
        """Generate two-column slide"""
        slide_layout = self.prs.slide_layouts[3] if len(self.prs.slide_layouts) > 3 else self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if isinstance(config.content, dict):
            # Add left column
            if 'left' in config.content and len(slide.placeholders) > 1:
                left_placeholder = slide.placeholders[1]
                self._add_content_to_placeholder(left_placeholder, config.content['left'])

            # Add right column
            if 'right' in config.content and len(slide.placeholders) > 2:
                right_placeholder = slide.placeholders[2]
                self._add_content_to_placeholder(right_placeholder, config.content['right'])

    def _add_content_to_placeholder(self, placeholder, content_items):
        """Add content items to a placeholder"""
        text_frame = placeholder.text_frame
        text_frame.clear()

        for item in content_items:
            if isinstance(item, dict):
                p = text_frame.add_paragraph()
                p.text = item.get('text', '')
                p.level = item.get('level', 0)
            elif isinstance(item, str):
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 0

    def _generate_chart_slide(self, config: SlideConfig):
        """Generate chart slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if config.content:
            for item in config.content:
                if isinstance(item, dict) and item.get('type') == 'chart':
                    self._add_chart(slide, item)

    def _add_chart(self, slide, chart_config: Dict):
        """Add chart to slide"""
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)

        chart_data = CategoryChartData()
        data = chart_config.get('data', {})

        # Add categories and values
        if 'categories' in data:
            chart_data.categories = data['categories']

        if 'series' in data:
            for series in data['series']:
                chart_data.add_series(series.get('name', ''), series.get('values', []))

        # Determine chart type
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }

        chart_type = chart_type_map.get(
            chart_config.get('chart_type', 'column'),
            XL_CHART_TYPE.COLUMN_CLUSTERED
        )

        # Add chart
        chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart

        # Apply chart options
        options = chart_config.get('options', {})
        if 'title' in options:
            chart.has_title = True
            chart.chart_title.text_frame.text = options['title']

    def _generate_table_slide(self, config: SlideConfig):
        """Generate table slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if config.content:
            for item in config.content:
                if isinstance(item, dict) and item.get('type') == 'table':
                    self._add_table(slide, item)

    def _add_table(self, slide, table_config: Dict):
        """Add table to slide"""
        headers = table_config.get('headers', [])
        rows = table_config.get('rows', [])

        if not headers or not rows:
            return

        rows_count = len(rows) + 1  # +1 for header
        cols_count = len(headers)

        x, y = Inches(1), Inches(2)
        cx, cy = Inches(8), Inches(0.5 * rows_count)

        table = slide.shapes.add_table(rows_count, cols_count, x, y, cx, cy).table

        # Set headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            # Style header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Add data
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data[:cols_count]):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)

    def _generate_code_slide(self, config: SlideConfig):
        """Generate code slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if config.content:
            for item in config.content:
                if isinstance(item, dict) and item.get('type') == 'code':
                    self._add_code_block(slide, item)

    def _add_code_block(self, slide, code_config: Dict):
        """Add code block to slide"""
        code = code_config.get('code', '')
        language = code_config.get('language', 'text')

        # Add text box for code
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Add code with monospace font
        p = text_frame.paragraphs[0]
        p.text = code
        p.font.name = 'Consolas'
        p.font.size = Pt(10)

        # Add background color
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(240, 240, 240)

    def _generate_quote_slide(self, config: SlideConfig):
        """Generate quote slide"""
        slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.content:
            for item in config.content:
                if isinstance(item, dict) and item.get('type') == 'quote':
                    # Add quote text
                    quote_box = slide.shapes.add_textbox(
                        Inches(1), Inches(2), Inches(8), Inches(2)
                    )
                    quote_text = quote_box.text_frame
                    p = quote_text.paragraphs[0]
                    p.text = f'"{item.get("text", "")}"'
                    p.font.size = Pt(24)
                    p.font.italic = True
                    p.alignment = PP_ALIGN.CENTER

                    # Add author
                    if item.get('author'):
                        author_box = slide.shapes.add_textbox(
                            Inches(1), Inches(4.5), Inches(8), Inches(0.5)
                        )
                        author_text = author_box.text_frame
                        p = author_text.paragraphs[0]
                        p.text = f"— {item['author']}"
                        p.font.size = Pt(18)
                        p.alignment = PP_ALIGN.CENTER

    def _generate_timeline_slide(self, config: SlideConfig):
        """Generate timeline slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if config.content:
            for item in config.content:
                if isinstance(item, dict) and item.get('type') == 'timeline':
                    self._add_timeline(slide, item)

    def _add_timeline(self, slide, timeline_config: Dict):
        """Add timeline to slide"""
        events = timeline_config.get('events', [])
        style = timeline_config.get('style', 'horizontal')

        if not events:
            return

        # Simple horizontal timeline
        line_y = Inches(3.5)
        line_left = Inches(1)
        line_width = Inches(8)

        # Add timeline line
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            line_left, line_y,
            line_left + line_width, line_y
        )
        line.line.width = Pt(3)

        # Add events
        event_spacing = line_width / (len(events) + 1)
        for i, event in enumerate(events):
            event_x = line_left + event_spacing * (i + 1)

            # Add event marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                event_x - Inches(0.15),
                line_y - Inches(0.15),
                Inches(0.3),
                Inches(0.3)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(255, 0, 0)

            # Add event text
            text_box = slide.shapes.add_textbox(
                event_x - Inches(0.75),
                line_y + Inches(0.3) if i % 2 == 0 else line_y - Inches(1),
                Inches(1.5),
                Inches(0.6)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]

            event_text = event.get('date', '')
            if event.get('title'):
                event_text += f"\\n{event['title']}"
            p.text = event_text
            p.font.size = Pt(10)

    def _generate_image_slide(self, config: SlideConfig):
        """Generate image slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            slide.shapes.title.text = config.title

        if config.content:
            for item in config.content:
                if isinstance(item, dict) and item.get('type') == 'image':
                    self._add_image(slide, item)

    def _add_image(self, slide, image_config: Dict):
        """Add image to slide"""
        src = image_config.get('src', '')

        if not src:
            return

        # Handle local files or URLs
        if src.startswith('http'):
            # Download image to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                urllib.request.urlretrieve(src, tmp.name)
                image_path = tmp.name
        else:
            image_path = src

        # Add image to slide
        if Path(image_path).exists():
            attributes = image_config.get('attributes', {})

            # Parse position and size
            left = Inches(float(attributes.get('x', 2.5)))
            top = Inches(float(attributes.get('y', 2)))
            width = Inches(float(attributes.get('width', 5)))

            slide.shapes.add_picture(image_path, left, top, width=width)

    def _generate_blank_slide(self, config: SlideConfig):
        """Generate blank slide"""
        slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if config.title:
            # Add title manually
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(8), Inches(1)
            )
            title_box.text = config.title
            title_box.text_frame.paragraphs[0].font.size = Pt(32)
            title_box.text_frame.paragraphs[0].font.bold = True

    def _apply_slide_settings(self, slide, config: SlideConfig):
        """Apply slide-specific settings"""

        # Background color
        if config.background:
            background = slide.background
            fill = background.fill
            fill.solid()

            # Parse color
            if config.background.startswith('#'):
                color = self._parse_color(config.background)
                fill.fore_color.rgb = color

        # Background image
        if config.metadata.get('bg_image'):
            # This would require more complex implementation
            pass

        # Speaker notes
        if config.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = config.notes

        # Transitions and animations would require XML manipulation
        # Store them as metadata for potential future implementation
        if config.transition:
            slide.transition_type = config.transition

        if config.animation:
            slide.animation_type = config.animation

class MarkdownToPowerPoint:
    """Main class for converting Markdown to PowerPoint"""

    def __init__(self):
        self.parser = MarkdownPresentationParser()
        self.generator = PowerPointGenerator()

    async def convert(self, markdown_content: str, output_path: str) -> Dict[str, Any]:
        """Convert markdown to PowerPoint presentation"""

        try:
            # Parse markdown
            config, slides = self.parser.parse(markdown_content)

            # Generate PowerPoint
            presentation = self.generator.generate(config, slides)

            # Save presentation
            presentation.save(output_path)

            return {
                'success': True,
                'output_path': output_path,
                'slide_count': len(slides),
                'title': config.title,
                'metadata': {
                    'author': config.author,
                    'theme': config.theme,
                    'aspect_ratio': config.aspect_ratio
                }
            }

        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

    async def convert_file(self, markdown_file: str, output_path: str) -> Dict[str, Any]:
        """Convert markdown file to PowerPoint"""

        with open(markdown_file, 'r', encoding='utf-8') as f:
            content = f.read()

        return await self.convert(content, output_path)

# Export classes for unified server
__all__ = [
    'SlideType',
    'SlideConfig',
    'PresentationConfig',
    'MarkdownPresentationParser',
    'PowerPointGenerator',
    'MarkdownToPowerPoint'
]
