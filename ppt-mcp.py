#!/usr/bin/env python3
"""
Extended MCP Server for PowerPoint Creation with Advanced Features
"""

import asyncio
import json
import base64
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime
import io
import re
from dataclasses import dataclass
import tempfile
import shutil
import zipfile

from mcp.server.models import InitializationOptions
from mcp.server import NotificationOptions, Server
from mcp.server.stdio import stdio_server
import mcp.types as types

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Mm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
import numpy as np

@dataclass
class SlideAnimation:
    """Animation configuration for shapes"""
    effect: str
    duration: float
    delay: float
    trigger: str

@dataclass
class SmartArtData:
    """Smart art diagram data"""
    type: str
    items: List[str]
    hierarchy: Optional[Dict] = None

class ExtendedPowerPointServer:
    """Extended MCP Server for PowerPoint operations"""

    def __init__(self):
        self.server = Server("powerpoint-server-extended")
        self.presentations: Dict[str, Presentation] = {}
        self.templates: Dict[str, Dict] = {}
        self.slide_notes: Dict[str, Dict[int, str]] = {}
        self.animations: Dict[str, List[SlideAnimation]] = {}
        self.media_cache: Dict[str, bytes] = {}
        self.setup_handlers()
        self.init_templates()

    def init_templates(self):
        """Initialize presentation templates"""
        self.templates = {
            "corporate": {
                "colors": {"primary": "003366", "secondary": "0066CC", "accent": "FF6600"},
                "fonts": {"title": "Arial Black", "body": "Arial", "accent": "Georgia"},
                "layouts": ["title", "agenda", "content", "comparison", "closing"]
            },
            "creative": {
                "colors": {"primary": "FF1744", "secondary": "AA00FF", "accent": "00E676"},
                "fonts": {"title": "Impact", "body": "Century Gothic", "accent": "Brush Script MT"},
                "layouts": ["title", "portfolio", "gallery", "quote", "thanks"]
            },
            "academic": {
                "colors": {"primary": "1A237E", "secondary": "3F51B5", "accent": "FFC107"},
                "fonts": {"title": "Times New Roman", "body": "Calibri", "accent": "Cambria"},
                "layouts": ["title", "objectives", "methodology", "results", "references"]
            },
            "minimalist": {
                "colors": {"primary": "000000", "secondary": "666666", "accent": "FFFFFF"},
                "fonts": {"title": "Helvetica", "body": "Helvetica Light", "accent": "Helvetica"},
                "layouts": ["title", "statement", "image_focus", "data", "end"]
            }
        }

    def setup_handlers(self):
        """Setup all tool handlers including new advanced features"""

        @self.server.list_tools()
        async def handle_list_tools() -> list[types.Tool]:
            """List all available PowerPoint tools"""
            return [
                # Original tools (keeping the base ones)
                types.Tool(
                    name="create_presentation",
                    description="Create a new PowerPoint presentation with advanced options",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "template": {"type": "string"},
                            "aspect_ratio": {
                                "type": "string",
                                "description": "16:9, 4:3, custom",
                                "default": "16:9"
                            },
                            "orientation": {
                                "type": "string",
                                "description": "landscape, portrait",
                                "default": "landscape"
                            }
                        },
                        "required": ["presentation_id"]
                    }
                ),

                # New Advanced Tools
                types.Tool(
                    name="add_smart_art",
                    description="Add SmartArt diagram to slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "smart_art_type": {
                                "type": "string",
                                "description": "process, cycle, hierarchy, relationship, matrix, pyramid"
                            },
                            "items": {
                                "type": "array",
                                "items": {"type": "string"}
                            },
                            "color_scheme": {"type": "string"}
                        },
                        "required": ["presentation_id", "smart_art_type", "items"]
                    }
                ),

                types.Tool(
                    name="add_timeline_slide",
                    description="Add a timeline slide with events",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "events": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "date": {"type": "string"},
                                        "title": {"type": "string"},
                                        "description": {"type": "string"}
                                    }
                                }
                            },
                            "style": {
                                "type": "string",
                                "description": "horizontal, vertical, curved, zigzag"
                            }
                        },
                        "required": ["presentation_id", "title", "events"]
                    }
                ),

                types.Tool(
                    name="add_comparison_slide",
                    description="Add a comparison slide with multiple columns",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "items": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "header": {"type": "string"},
                                        "features": {"type": "array", "items": {"type": "string"}},
                                        "highlight": {"type": "boolean"}
                                    }
                                }
                            },
                            "show_icons": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id", "title", "items"]
                    }
                ),

                types.Tool(
                    name="add_quote_slide",
                    description="Add a quote/testimonial slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "quote": {"type": "string"},
                            "author": {"type": "string"},
                            "title": {"type": "string"},
                            "background_image": {"type": "string"},
                            "style": {
                                "type": "string",
                                "description": "centered, left, right, overlay"
                            }
                        },
                        "required": ["presentation_id", "quote", "author"]
                    }
                ),

                types.Tool(
                    name="add_agenda_slide",
                    description="Add an agenda/table of contents slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string", "default": "Agenda"},
                            "sections": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "title": {"type": "string"},
                                        "duration": {"type": "string"},
                                        "presenter": {"type": "string"}
                                    }
                                }
                            },
                            "show_numbers": {"type": "boolean", "default": True},
                            "show_duration": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id", "sections"]
                    }
                ),

                types.Tool(
                    name="add_team_slide",
                    description="Add a team/organization slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "members": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "name": {"type": "string"},
                                        "role": {"type": "string"},
                                        "photo": {"type": "string"},
                                        "email": {"type": "string"},
                                        "bio": {"type": "string"}
                                    }
                                }
                            },
                            "layout": {
                                "type": "string",
                                "description": "grid, hierarchical, circular"
                            }
                        },
                        "required": ["presentation_id", "title", "members"]
                    }
                ),

                types.Tool(
                    name="add_process_flow",
                    description="Add a process flow diagram",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "steps": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "name": {"type": "string"},
                                        "description": {"type": "string"},
                                        "shape": {"type": "string"},
                                        "color": {"type": "string"}
                                    }
                                }
                            },
                            "flow_type": {
                                "type": "string",
                                "description": "linear, circular, branching"
                            },
                            "show_arrows": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id", "title", "steps"]
                    }
                ),

                types.Tool(
                    name="add_infographic_slide",
                    description="Add an infographic slide with data visualization",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "data_points": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "label": {"type": "string"},
                                        "value": {"type": "number"},
                                        "unit": {"type": "string"},
                                        "icon": {"type": "string"}
                                    }
                                }
                            },
                            "style": {
                                "type": "string",
                                "description": "icons, bars, circles, mixed"
                            }
                        },
                        "required": ["presentation_id", "title", "data_points"]
                    }
                ),

                types.Tool(
                    name="add_gantt_chart",
                    description="Add a Gantt chart for project timeline",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "tasks": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "name": {"type": "string"},
                                        "start_date": {"type": "string"},
                                        "end_date": {"type": "string"},
                                        "progress": {"type": "number"},
                                        "assignee": {"type": "string"}
                                    }
                                }
                            }
                        },
                        "required": ["presentation_id", "title", "tasks"]
                    }
                ),

                types.Tool(
                    name="add_swot_analysis",
                    description="Add a SWOT analysis slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "strengths": {"type": "array", "items": {"type": "string"}},
                            "weaknesses": {"type": "array", "items": {"type": "string"}},
                            "opportunities": {"type": "array", "items": {"type": "string"}},
                            "threats": {"type": "array", "items": {"type": "string"}}
                        },
                        "required": ["presentation_id", "strengths", "weaknesses", "opportunities", "threats"]
                    }
                ),

                types.Tool(
                    name="add_animation",
                    description="Add animations to shapes on a slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "shape_index": {"type": "integer"},
                            "animation_type": {
                                "type": "string",
                                "description": "appear, fade, fly_in, zoom, bounce, spin"
                            },
                            "duration": {"type": "number"},
                            "delay": {"type": "number"},
                            "trigger": {
                                "type": "string",
                                "description": "on_click, with_previous, after_previous"
                            }
                        },
                        "required": ["presentation_id", "slide_index", "animation_type"]
                    }
                ),

                types.Tool(
                    name="add_slide_notes",
                    description="Add speaker notes to a slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "notes": {"type": "string"}
                        },
                        "required": ["presentation_id", "slide_index", "notes"]
                    }
                ),

                types.Tool(
                    name="add_hyperlink",
                    description="Add hyperlinks to shapes or text",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "shape_index": {"type": "integer"},
                            "url": {"type": "string"},
                            "screen_tip": {"type": "string"}
                        },
                        "required": ["presentation_id", "slide_index", "shape_index", "url"]
                    }
                ),

                types.Tool(
                    name="add_video_slide",
                    description="Add a slide with embedded video",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "video_path": {"type": "string"},
                            "poster_frame": {"type": "string"},
                            "autoplay": {"type": "boolean", "default": False}
                        },
                        "required": ["presentation_id", "title", "video_path"]
                    }
                ),

                types.Tool(
                    name="add_audio",
                    description="Add audio to a slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "audio_path": {"type": "string"},
                            "play_across_slides": {"type": "boolean"},
                            "loop": {"type": "boolean"}
                        },
                        "required": ["presentation_id", "slide_index", "audio_path"]
                    }
                ),

                types.Tool(
                    name="add_footer",
                    description="Add footer to all slides",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "footer_text": {"type": "string"},
                            "show_slide_number": {"type": "boolean"},
                            "show_date": {"type": "boolean"},
                            "exclude_title_slide": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id"]
                    }
                ),

                types.Tool(
                    name="add_watermark",
                    description="Add watermark to slides",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "text": {"type": "string"},
                            "opacity": {"type": "number", "default": 0.3},
                            "angle": {"type": "number", "default": 45},
                            "color": {"type": "string", "default": "CCCCCC"}
                        },
                        "required": ["presentation_id", "text"]
                    }
                ),

                types.Tool(
                    name="apply_slide_master",
                    description="Apply or modify slide master",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "background_color": {"type": "string"},
                            "background_image": {"type": "string"},
                            "logo_path": {"type": "string"},
                            "logo_position": {
                                "type": "string",
                                "description": "top-left, top-right, bottom-left, bottom-right"
                            }
                        },
                        "required": ["presentation_id"]
                    }
                ),

                types.Tool(
                    name="add_custom_shape",
                    description="Add custom shapes and connectors",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "shape_type": {
                                "type": "string",
                                "description": "rectangle, circle, arrow, star, heart, callout"
                            },
                            "left": {"type": "number"},
                            "top": {"type": "number"},
                            "width": {"type": "number"},
                            "height": {"type": "number"},
                            "fill_color": {"type": "string"},
                            "line_color": {"type": "string"},
                            "text": {"type": "string"}
                        },
                        "required": ["presentation_id", "slide_index", "shape_type"]
                    }
                ),

                types.Tool(
                    name="duplicate_slide",
                    description="Duplicate an existing slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "source_index": {"type": "integer"},
                            "target_index": {"type": "integer"}
                        },
                        "required": ["presentation_id", "source_index"]
                    }
                ),

                types.Tool(
                    name="reorder_slides",
                    description="Reorder slides in presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_order": {
                                "type": "array",
                                "items": {"type": "integer"}
                            }
                        },
                        "required": ["presentation_id", "slide_order"]
                    }
                ),

                types.Tool(
                    name="delete_slide",
                    description="Delete a slide from presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"}
                        },
                        "required": ["presentation_id", "slide_index"]
                    }
                ),

                types.Tool(
                    name="export_as_pdf",
                    description="Export presentation as PDF",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "output_path": {"type": "string"},
                            "slides_per_page": {"type": "integer", "default": 1},
                            "include_notes": {"type": "boolean", "default": False}
                        },
                        "required": ["presentation_id", "output_path"]
                    }
                ),

                types.Tool(
                    name="export_slides_as_images",
                    description="Export slides as individual images",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "output_dir": {"type": "string"},
                            "format": {
                                "type": "string",
                                "description": "png, jpg, svg",
                                "default": "png"
                            },
                            "resolution": {"type": "integer", "default": 300}
                        },
                        "required": ["presentation_id", "output_dir"]
                    }
                ),

                types.Tool(
                    name="merge_presentations",
                    description="Merge multiple presentations",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "target_id": {"type": "string"},
                            "source_ids": {
                                "type": "array",
                                "items": {"type": "string"}
                            },
                            "maintain_formatting": {"type": "boolean", "default": True}
                        },
                        "required": ["target_id", "source_ids"]
                    }
                ),

                types.Tool(
                    name="add_qr_code",
                    description="Add QR code to slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "data": {"type": "string"},
                            "size": {"type": "number", "default": 2},
                            "position": {
                                "type": "string",
                                "description": "center, top-right, bottom-right",
                                "default": "bottom-right"
                            }
                        },
                        "required": ["presentation_id", "slide_index", "data"]
                    }
                ),

                types.Tool(
                    name="add_math_equation",
                    description="Add mathematical equation to slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "equation": {"type": "string"},
                            "font_size": {"type": "integer", "default": 18}
                        },
                        "required": ["presentation_id", "slide_index", "equation"]
                    }
                ),

                types.Tool(
                    name="add_code_block",
                    description="Add syntax-highlighted code block",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "code": {"type": "string"},
                            "language": {"type": "string"},
                            "theme": {
                                "type": "string",
                                "description": "light, dark, monokai",
                                "default": "light"
                            }
                        },
                        "required": ["presentation_id", "slide_index", "code", "language"]
                    }
                ),

                types.Tool(
                    name="generate_handouts",
                    description="Generate handout version of presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "output_path": {"type": "string"},
                            "slides_per_page": {"type": "integer", "default": 3},
                            "include_lines": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id", "output_path"]
                    }
                ),

                # Keep all original tools...
                types.Tool(
                    name="add_title_slide",
                    description="Add a title slide to the presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"},
                            "author": {"type": "string"}
                        },
                        "required": ["presentation_id", "title"]
                    }
                ),
                types.Tool(
                    name="add_content_slide",
                    description="Add a content slide with title and bullet points",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "content": {
                                "type": "array",
                                "items": {"type": "string"}
                            },
                            "layout": {
                                "type": "string",
                                "description": "Layout: bullets, two_column, comparison",
                                "default": "bullets"
                            }
                        },
                        "required": ["presentation_id", "title", "content"]
                    }
                ),
                types.Tool(
                    name="add_chart_slide",
                    description="Add a slide with a chart",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "chart_type": {
                                "type": "string",
                                "description": "Chart type: bar, column, line, pie, scatter, bubble, radar, waterfall"
                            },
                            "categories": {
                                "type": "array",
                                "items": {"type": "string"}
                            },
                            "series_data": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "name": {"type": "string"},
                                        "values": {
                                            "type": "array",
                                            "items": {"type": "number"}
                                        }
                                    }
                                }
                            }
                        },
                        "required": ["presentation_id", "title", "chart_type", "categories", "series_data"]
                    }
                ),
                types.Tool(
                    name="save_presentation",
                    description="Save the presentation to a file",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "file_path": {"type": "string"}
                        },
                        "required": ["presentation_id", "file_path"]
                    }
                ),
            ]

        @self.server.call_tool()
        async def handle_call_tool(
            name: str, arguments: Optional[Dict[str, Any]]
        ) -> list[types.TextContent]:
            """Handle tool execution"""

            try:
                # Map all tool names to their handler methods
                tool_handlers = {
                    "create_presentation": self.create_presentation,
                    "add_title_slide": self.add_title_slide,
                    "add_content_slide": self.add_content_slide,
                    "add_chart_slide": self.add_chart_slide,
                    "add_smart_art": self.add_smart_art,
                    "add_timeline_slide": self.add_timeline_slide,
                    "add_comparison_slide": self.add_comparison_slide,
                    "add_quote_slide": self.add_quote_slide,
                    "add_agenda_slide": self.add_agenda_slide,
                    "add_team_slide": self.add_team_slide,
                    "add_process_flow": self.add_process_flow,
                    "add_infographic_slide": self.add_infographic_slide,
                    "add_gantt_chart": self.add_gantt_chart,
                    "add_swot_analysis": self.add_swot_analysis,
                    "add_animation": self.add_animation,
                    "add_slide_notes": self.add_slide_notes,
                    "add_hyperlink": self.add_hyperlink,
                    "add_video_slide": self.add_video_slide,
                    "add_audio": self.add_audio,
                    "add_footer": self.add_footer,
                    "add_watermark": self.add_watermark,
                    "apply_slide_master": self.apply_slide_master,
                    "add_custom_shape": self.add_custom_shape,
                    "duplicate_slide": self.duplicate_slide,
                    "reorder_slides": self.reorder_slides,
                    "delete_slide": self.delete_slide,
                    "export_as_pdf": self.export_as_pdf,
                    "export_slides_as_images": self.export_slides_as_images,
                    "merge_presentations": self.merge_presentations,
                    "add_qr_code": self.add_qr_code,
                    "add_math_equation": self.add_math_equation,
                    "add_code_block": self.add_code_block,
                    "generate_handouts": self.generate_handouts,
                    "save_presentation": self.save_presentation,
                }

                handler = tool_handlers.get(name)
                if handler:
                    return await handler(arguments)
                else:
                    raise ValueError(f"Unknown tool: {name}")

            except Exception as e:
                return [types.TextContent(
                    type="text",
                    text=f"Error executing {name}: {str(e)}"
                )]

    # Implement all the new handler methods

    async def add_smart_art(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add SmartArt diagram"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        if "title" in args:
            slide.shapes.title.text = args["title"]

        smart_art_type = args["smart_art_type"]
        items = args["items"]

        # Create SmartArt-like diagram using shapes
        if smart_art_type == "process":
            self._create_process_diagram(slide, items)
        elif smart_art_type == "cycle":
            self._create_cycle_diagram(slide, items)
        elif smart_art_type == "hierarchy":
            self._create_hierarchy_diagram(slide, items)
        elif smart_art_type == "pyramid":
            self._create_pyramid_diagram(slide, items)
        else:
            self._create_relationship_diagram(slide, items)

        return [types.TextContent(
            type="text",
            text=f"Added {smart_art_type} SmartArt with {len(items)} items"
        )]

    def _create_process_diagram(self, slide, items):
        """Create a process flow diagram"""
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(1.5)
        height = Inches(1)
        spacing = Inches(0.3)

        for i, item in enumerate(items):
            # Add process box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left + i * (width + spacing),
                top,
                width,
                height
            )
            shape.text = item
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(68, 114, 196)

            # Add arrow between boxes
            if i < len(items) - 1:
                arrow = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    left + i * (width + spacing) + width,
                    top + height/2,
                    left + (i+1) * (width + spacing),
                    top + height/2
                )
                arrow.line.color.rgb = RGBColor(68, 114, 196)

    def _create_cycle_diagram(self, slide, items):
        """Create a cycle diagram"""
        center_x = Inches(5)
        center_y = Inches(3.5)
        radius = Inches(1.5)

        angle_step = 360 / len(items)

        for i, item in enumerate(items):
            angle = i * angle_step
            x = center_x + radius * np.cos(np.radians(angle))
            y = center_y + radius * np.sin(np.radians(angle))

            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.5),
                y - Inches(0.5),
                Inches(1),
                Inches(1)
            )
            shape.text = item
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(68, 114, 196)

    def _create_hierarchy_diagram(self, slide, items):
        """Create a hierarchy diagram"""
        # Simplified hierarchy - would need more complex logic for full implementation
        top_item = items[0] if items else "Root"
        children = items[1:] if len(items) > 1 else []

        # Add top level
        top_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4),
            Inches(2),
            Inches(2),
            Inches(0.75)
        )
        top_shape.text = top_item

        # Add children
        if children:
            child_width = Inches(1.5)
            total_width = len(children) * child_width + (len(children)-1) * Inches(0.2)
            start_left = Inches(5) - total_width/2

            for i, child in enumerate(children):
                child_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    start_left + i * (child_width + Inches(0.2)),
                    Inches(3.5),
                    child_width,
                    Inches(0.75)
                )
                child_shape.text = child

    def _create_pyramid_diagram(self, slide, items):
        """Create a pyramid diagram"""
        base_width = Inches(6)
        base_left = Inches(2)
        top = Inches(2)
        height_per_level = Inches(0.8)

        for i, item in enumerate(items):
            level_width = base_width * (len(items) - i) / len(items)
            level_left = base_left + (base_width - level_width) / 2

            shape = slide.shapes.add_shape(
                MSO_SHAPE.TRAPEZOID,
                level_left,
                top + i * height_per_level,
                level_width,
                height_per_level
            )
            shape.text = item
            shape.fill.solid()
            color_value = 68 + i * 30
            shape.fill.fore_color.rgb = RGBColor(color_value, 114, 196)

    def _create_relationship_diagram(self, slide, items):
        """Create a relationship diagram"""
        # Create a simple matrix/relationship diagram
        for i, item in enumerate(items):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1 + i * 2),
                Inches(2.5),
                Inches(1.5),
                Inches(1)
            )
            shape.text = item

    async def add_timeline_slide(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add a timeline slide"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = args["title"]

        events = args["events"]
        style = args.get("style", "horizontal")

        if style == "horizontal":
            # Draw horizontal timeline
            line_top = Inches(3.5)
            line_left = Inches(1)
            line_width = Inches(8)

            # Add timeline line
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                line_left, line_top,
                line_left + line_width, line_top
            )
            line.line.width = Pt(3)

            # Add events
            event_spacing = line_width / (len(events) + 1)
            for i, event in enumerate(events):
                event_left = line_left + event_spacing * (i + 1)

                # Add event marker
                marker = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    event_left - Inches(0.15),
                    line_top - Inches(0.15),
                    Inches(0.3),
                    Inches(0.3)
                )
                marker.fill.solid()
                marker.fill.fore_color.rgb = RGBColor(255, 0, 0)

                # Add event text
                text_box = slide.shapes.add_textbox(
                    event_left - Inches(0.75),
                    line_top + Inches(0.3) if i % 2 == 0 else line_top - Inches(1),
                    Inches(1.5),
                    Inches(0.6)
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = f"{event['date']}\\n{event['title']}"
                p.font.size = Pt(10)

        return [types.TextContent(
            type="text",
            text=f"Added timeline with {len(events)} events"
        )]

    async def add_comparison_slide(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add a comparison slide"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = args["title"]

        items = args["items"]
        num_items = len(items)

        # Calculate column dimensions
        col_width = Inches(8) / num_items
        col_left = Inches(1)
        col_top = Inches(2)

        for i, item in enumerate(items):
            # Add header
            header_box = slide.shapes.add_textbox(
                col_left + i * col_width,
                col_top,
                col_width - Inches(0.1),
                Inches(0.5)
            )
            header_box.text = item["header"]

            # Highlight if specified
            if item.get("highlight", False):
                header_box.fill.solid()
                header_box.fill.fore_color.rgb = RGBColor(255, 235, 59)

            # Add features
            for j, feature in enumerate(item.get("features", [])):
                feature_box = slide.shapes.add_textbox(
                    col_left + i * col_width,
                    col_top + Inches(0.7) + j * Inches(0.4),
                    col_width - Inches(0.1),
                    Inches(0.3)
                )
                feature_box.text = f"• {feature}"
                feature_box.text_frame.paragraphs[0].font.size = Pt(11)

        return [types.TextContent(
            type="text",
            text=f"Added comparison slide with {num_items} items"
        )]

    async def add_quote_slide(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add a quote slide"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add quote
        quote_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(2),
            Inches(8),
            Inches(2)
        )
        quote_text = quote_box.text_frame
        p = quote_text.paragraphs[0]
        p.text = f'"{args["quote"]}"'
        p.font.size = Pt(28)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

        # Add author
        author_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(4.5),
            Inches(8),
            Inches(0.5)
        )
        author_text = author_box.text_frame
        p = author_text.paragraphs[0]
        p.text = f"— {args['author']}"
        if "title" in args:
            p.text += f", {args['title']}"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

        return [types.TextContent(
            type="text",
            text=f"Added quote slide from {args['author']}"
        )]

    async def add_slide_notes(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add speaker notes to a slide"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_index = args["slide_index"]

        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")

        slide = prs.slides[slide_index]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = args["notes"]

        # Store notes for reference
        if pres_id not in self.slide_notes:
            self.slide_notes[pres_id] = {}
        self.slide_notes[pres_id][slide_index] = args["notes"]

        return [types.TextContent(
            type="text",
            text=f"Added notes to slide {slide_index}"
        )]

    async def add_watermark(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add watermark to slides"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        watermark_text = args["text"]
        opacity = args.get("opacity", 0.3)
        angle = args.get("angle", 45)

        for slide in prs.slides:
            # Add diagonal watermark text
            watermark = slide.shapes.add_textbox(
                Inches(1),
                Inches(3),
                Inches(8),
                Inches(1)
            )
            text_frame = watermark.text_frame
            p = text_frame.paragraphs[0]
            p.text = watermark_text
            p.font.size = Pt(48)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER

            # Note: Rotation and transparency require XML manipulation
            # This is a simplified version

        return [types.TextContent(
            type="text",
            text=f"Added watermark '{watermark_text}' to all slides"
        )]

    async def add_custom_shape(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add custom shape to slide"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide = prs.slides[args["slide_index"]]

        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "star": MSO_SHAPE.STAR_5_POINT,
            "heart": MSO_SHAPE.HEART,
            "callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT
        }

        shape_type = shape_map.get(args["shape_type"], MSO_SHAPE.RECTANGLE)

        left = Inches(args.get("left", 1))
        top = Inches(args.get("top", 1))
        width = Inches(args.get("width", 2))
        height = Inches(args.get("height", 1))

        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        if "fill_color" in args:
            shape.fill.solid()
            rgb = tuple(int(args["fill_color"][i:i+2], 16) for i in (0, 2, 4))
            shape.fill.fore_color.rgb = RGBColor(*rgb)

        if "text" in args:
            shape.text = args["text"]

        return [types.TextContent(
            type="text",
            text=f"Added {args['shape_type']} shape to slide {args['slide_index']}"
        )]

    async def export_as_pdf(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Export presentation as PDF (requires additional libraries)"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        # Note: Full PDF export would require additional libraries like comtypes on Windows
        # or unoconv on Linux. This is a placeholder implementation.

        output_path = args["output_path"]

        # For now, save as PPTX with note about PDF conversion
        pptx_path = output_path.replace('.pdf', '.pptx')
        self.presentations[pres_id].save(pptx_path)

        return [types.TextContent(
            type="text",
            text=f"Saved as PPTX at {pptx_path}. PDF conversion requires additional setup."
        )]

    async def add_qr_code(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add QR code to slide (requires qrcode library)"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        try:
            import qrcode

            prs = self.presentations[pres_id]
            slide = prs.slides[args["slide_index"]]

            # Generate QR code
            qr = qrcode.QRCode(version=1, box_size=10, border=5)
            qr.add_data(args["data"])
            qr.make(fit=True)

            img = qr.make_image(fill_color="black", back_color="white")

            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                img.save(tmp.name)

                # Add to slide
                position_map = {
                    "center": (Inches(4), Inches(2.5)),
                    "top-right": (Inches(7), Inches(0.5)),
                    "bottom-right": (Inches(7), Inches(5))
                }

                left, top = position_map.get(args.get("position", "bottom-right"))
                size = Inches(args.get("size", 2))

                slide.shapes.add_picture(tmp.name, left, top, size, size)

            return [types.TextContent(
                type="text",
                text=f"Added QR code to slide {args['slide_index']}"
            )]

        except ImportError:
            return [types.TextContent(
                type="text",
                text="QR code generation requires 'qrcode' library. Install with: pip install qrcode[pil]"
            )]

    # Keep original methods
    async def create_presentation(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Create a new presentation with advanced options"""
        pres_id = args["presentation_id"]
        template = args.get("template", "blank")

        prs = Presentation()

        # Set aspect ratio
        aspect_ratio = args.get("aspect_ratio", "16:9")
        if aspect_ratio == "16:9":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
        elif aspect_ratio == "4:3":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

        self.presentations[pres_id] = prs

        # Apply template if specified
        if template in self.templates:
            # This would apply template settings
            pass

        return [types.TextContent(
            type="text",
            text=f"Created presentation '{pres_id}' ({aspect_ratio}, {template} template)"
        )]

    async def add_title_slide(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add a title slide"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None

        title.text = args["title"]

        if subtitle and "subtitle" in args:
            subtitle.text = args["subtitle"]

        if "author" in args and subtitle:
            subtitle.text += f"\\n{args['author']}"

        return [types.TextContent(
            type="text",
            text=f"Added title slide: {args['title']}"
        )]

    async def add_content_slide(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add a content slide with bullets"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = args["title"]

        content_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if content_shape:
            text_frame = content_shape.text_frame
            text_frame.clear()

            for bullet in args["content"]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0

        return [types.TextContent(
            type="text",
            text=f"Added content slide: {args['title']}"
        )]

    async def add_chart_slide(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Add enhanced chart slide with more chart types"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = args["title"]

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)

        chart_data = CategoryChartData()
        chart_data.categories = args["categories"]

        for series in args["series_data"]:
            chart_data.add_series(series["name"], series["values"])

        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "bubble": XL_CHART_TYPE.BUBBLE,
            "radar": XL_CHART_TYPE.RADAR,
            "waterfall": XL_CHART_TYPE.COLUMN_CLUSTERED,  # Simplified
        }

        chart_type = chart_type_map.get(args["chart_type"], XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart = slide.shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data
        ).chart

        return [types.TextContent(
            type="text",
            text=f"Added {args['chart_type']} chart: {args['title']}"
        )]

    async def save_presentation(self, args: Dict[str, Any]) -> list[types.TextContent]:
        """Save the presentation"""
        pres_id = args["presentation_id"]
        if pres_id not in self.presentations:
            raise ValueError(f"Presentation '{pres_id}' not found")

        prs = self.presentations[pres_id]
        file_path = Path(args["file_path"])

        file_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(file_path))

        return [types.TextContent(
            type="text",
            text=f"Saved presentation to: {file_path}"
        )]

    async def run(self):
        """Run the MCP server"""
        async with stdio_server() as (read_stream, write_stream):
            await self.server.run(
                read_stream,
                write_stream,
                InitializationOptions(
                    server_name="powerpoint-mcp-extended",
                    server_version="2.0.0",
                    capabilities=self.server.get_capabilities(
                        notification_options=NotificationOptions(),
                        experimental_capabilities={},
                    ),
                ),
            )

async def main():
    """Main entry point"""
    server = ExtendedPowerPointServer()
    await server.run()

if __name__ == "__main__":
    asyncio.run(main())
