#!/usr/bin/env python3
"""
Unified PowerPoint MCP Server
Integrates: Markdown Conversion, MCP Tools, and Material Design
"""

import asyncio
import json
import base64
import re
import yaml
import colorsys
import random
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime
from dataclasses import dataclass, field
from enum import Enum
import tempfile
import hashlib
import urllib.request
from io import BytesIO

from mcp.server.models import InitializationOptions
from mcp.server import NotificationOptions, Server
from mcp.server.stdio import stdio_server
import mcp.types as types

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Mm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR

try:
    from PIL import Image
    import numpy as np
    import markdown
    import qrcode
except ImportError as e:
    print(f"Warning: Optional dependency missing: {e}")


# ============================================================================
# DATA CLASSES AND ENUMS
# ============================================================================

class SlideType(Enum):
    """Supported slide types"""
    TITLE = "title"
    SECTION = "section"
    CONTENT = "content"
    TWO_COLUMN = "two_column"
    IMAGE = "image"
    CHART = "chart"
    TABLE = "table"
    QUOTE = "quote"
    CODE = "code"
    TIMELINE = "timeline"
    COMPARISON = "comparison"
    TEAM = "team"
    BLANK = "blank"


@dataclass
class SlideConfig:
    """Configuration for a single slide"""
    type: SlideType
    title: Optional[str] = None
    subtitle: Optional[str] = None
    content: Optional[List[str]] = None
    image_url: Optional[str] = None
    chart_data: Optional[Dict] = None
    table_data: Optional[List[List[str]]] = None
    layout: Optional[str] = None
    background: Optional[str] = None
    notes: Optional[str] = None
    animation: Optional[str] = None
    theme_override: Optional[Dict] = None
    metadata: Dict = field(default_factory=dict)


@dataclass
class PresentationConfig:
    """Global presentation configuration"""
    title: str
    author: Optional[str] = None
    theme: str = "default"
    aspect_ratio: str = "16:9"
    footer: Optional[str] = None
    slide_numbers: bool = True
    date: bool = False
    company: Optional[str] = None
    logo: Optional[str] = None
    color_scheme: Optional[Dict] = None
    font_family: Optional[str] = None
    metadata: Dict = field(default_factory=dict)


@dataclass
class MaterialTheme:
    """Material Design Theme Configuration"""
    name: str
    primary_color: str
    primary_variant: str
    secondary_color: str
    secondary_variant: str
    background: str
    surface: str
    error: str
    on_primary: str
    on_secondary: str
    on_background: str
    on_surface: str
    on_error: str
    typography: Dict[str, Dict]
    elevation_shadows: List[Dict]
    spacing: Dict[str, float]
    corner_radius: float


# ============================================================================
# MATERIAL DESIGN COLOR PALETTES
# ============================================================================

class MaterialColors:
    """Material Design 3 Color Palettes"""
    
    RED = {"50": "FFEBEE", "500": "F44336", "900": "B71C1C"}
    PINK = {"50": "FCE4EC", "500": "E91E63", "900": "880E4F"}
    PURPLE = {"50": "F3E5F5", "500": "9C27B0", "900": "4A148C"}
    BLUE = {"50": "E3F2FD", "500": "2196F3", "900": "0D47A1"}
    GREEN = {"50": "E8F5E9", "500": "4CAF50", "900": "1B5E20"}
    ORANGE = {"50": "FFF3E0", "500": "FF9800", "900": "E65100"}
    GREY = {"50": "FAFAFA", "500": "9E9E9E", "900": "212121"}


# ============================================================================
# UNIFIED POWERPOINT SERVER
# ============================================================================

class UnifiedPowerPointServer:
    """Unified MCP Server for all PowerPoint operations"""

    def __init__(self):
        self.server = Server("unified-powerpoint-server")
        self.presentations: Dict[str, Presentation] = {}
        self.templates: Dict[str, Dict] = {}
        self.slide_notes: Dict[str, Dict[int, str]] = {}
        self.media_cache: Dict[str, bytes] = {}
        self.material_themes = self._init_material_themes()
        self.setup_handlers()
        self.init_templates()

    def init_templates(self):
        """Initialize presentation templates"""
        self.templates = {
            "corporate": {
                "colors": {"primary": "003366", "secondary": "0066CC", "accent": "FF6600"},
                "fonts": {"title": "Arial Black", "body": "Arial"},
                "layouts": ["title", "agenda", "content", "comparison", "closing"]
            },
            "creative": {
                "colors": {"primary": "FF1744", "secondary": "AA00FF", "accent": "00E676"},
                "fonts": {"title": "Impact", "body": "Century Gothic"},
                "layouts": ["title", "portfolio", "gallery", "quote"]
            },
            "academic": {
                "colors": {"primary": "1A237E", "secondary": "3F51B5", "accent": "FFC107"},
                "fonts": {"title": "Times New Roman", "body": "Calibri"},
                "layouts": ["title", "objectives", "methodology", "results"]
            },
            "minimalist": {
                "colors": {"primary": "000000", "secondary": "666666", "accent": "FFFFFF"},
                "fonts": {"title": "Helvetica", "body": "Helvetica Light"},
                "layouts": ["title", "statement", "image_focus", "data"]
            }
        }

    def _init_material_themes(self) -> Dict[str, MaterialTheme]:
        """Initialize Material Design themes"""
        return {
            "material_baseline": MaterialTheme(
                name="Material Baseline",
                primary_color="6200EE", primary_variant="3700B3",
                secondary_color="03DAC6", secondary_variant="018786",
                background="FFFFFF", surface="FFFFFF", error="B00020",
                on_primary="FFFFFF", on_secondary="000000",
                on_background="000000", on_surface="000000", on_error="FFFFFF",
                typography=self._get_typography(), elevation_shadows=self._get_shadows(),
                spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
                corner_radius=0.25
            ),
            "google_blue": MaterialTheme(
                name="Google Blue",
                primary_color="4285F4", primary_variant="1967D2",
                secondary_color="EA4335", secondary_variant="C5221F",
                background="FFFFFF", surface="F8F9FA", error="EA4335",
                on_primary="FFFFFF", on_secondary="FFFFFF",
                on_background="202124", on_surface="202124", on_error="FFFFFF",
                typography=self._get_typography(), elevation_shadows=self._get_shadows(),
                spacing={"xs": 0.25, "sm": 0.5, "md": 1.0, "lg": 1.5, "xl": 2.0},
                corner_radius=0.25
            )
        }

    @staticmethod
    def _get_typography() -> Dict:
        """Material Design typography scale"""
        return {
            "h1": {"size": 96, "weight": "light", "spacing": -1.5},
            "h2": {"size": 60, "weight": "light", "spacing": -0.5},
            "h3": {"size": 48, "weight": "regular", "spacing": 0},
            "h4": {"size": 34, "weight": "regular", "spacing": 0.25},
            "h5": {"size": 24, "weight": "regular", "spacing": 0},
            "h6": {"size": 20, "weight": "medium", "spacing": 0.15},
            "body1": {"size": 16, "weight": "regular", "spacing": 0.5},
            "body2": {"size": 14, "weight": "regular", "spacing": 0.25}
        }

    @staticmethod
    def _get_shadows() -> List[Dict]:
        """Material Design elevation shadows"""
        return [
            {"elevation": 0, "shadow": None},
            {"elevation": 1, "shadow": "0px 2px 1px -1px rgba(0,0,0,0.2)"},
            {"elevation": 2, "shadow": "0px 3px 1px -2px rgba(0,0,0,0.2)"},
            {"elevation": 4, "shadow": "0px 2px 4px -1px rgba(0,0,0,0.2)"}
        ]

    def setup_handlers(self):
        """Setup all MCP tool handlers"""

        @self.server.list_tools()
        async def handle_list_tools() -> list[types.Tool]:
            """List all available PowerPoint tools"""
            return [
                # === PRESENTATION MANAGEMENT ===
                types.Tool(
                    name="create_presentation",
                    description="Create a new PowerPoint presentation with optional template",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string", "description": "Unique ID for the presentation"},
                            "title": {"type": "string", "description": "Presentation title"},
                            "template": {"type": "string", "enum": ["corporate", "creative", "academic", "minimalist"], "description": "Template style"},
                            "aspect_ratio": {"type": "string", "enum": ["16:9", "4:3"], "default": "16:9"}
                        },
                        "required": ["presentation_id"]
                    }
                ),
                
                # === MARKDOWN CONVERSION ===
                types.Tool(
                    name="convert_markdown_to_pptx",
                    description="Convert markdown content directly to PowerPoint presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "markdown_content": {"type": "string", "description": "Markdown content with YAML frontmatter"},
                            "output_path": {"type": "string", "description": "Path to save the .pptx file"},
                            "presentation_id": {"type": "string", "description": "Optional ID to store in server"}
                        },
                        "required": ["markdown_content", "output_path"]
                    }
                ),
                
                types.Tool(
                    name="convert_markdown_file_to_pptx",
                    description="Convert markdown file to PowerPoint presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "markdown_file": {"type": "string", "description": "Path to markdown file"},
                            "output_path": {"type": "string", "description": "Path to save the .pptx file"}
                        },
                        "required": ["markdown_file", "output_path"]
                    }
                ),
                
                # === SLIDE CREATION ===
                types.Tool(
                    name="add_title_slide",
                    description="Add a title slide to the presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"}
                        },
                        "required": ["presentation_id", "title"]
                    }
                ),
                
                types.Tool(
                    name="add_content_slide",
                    description="Add a bullet-point content slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "content": {"type": "array", "items": {"type": "string"}, "description": "List of bullet points"}
                        },
                        "required": ["presentation_id", "title", "content"]
                    }
                ),
                
                types.Tool(
                    name="add_two_column_slide",
                    description="Add a two-column layout slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "left_content": {"type": "array", "items": {"type": "string"}},
                            "right_content": {"type": "array", "items": {"type": "string"}}
                        },
                        "required": ["presentation_id", "title", "left_content", "right_content"]
                    }
                ),
                
                # === CHARTS ===
                types.Tool(
                    name="add_chart_slide",
                    description="Add a chart slide (column, bar, line, or pie chart)",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "chart_type": {"type": "string", "enum": ["column", "bar", "line", "pie"]},
                            "categories": {"type": "array", "items": {"type": "string"}},
                            "series": {"type": "array", "items": {
                                "type": "object",
                                "properties": {
                                    "name": {"type": "string"},
                                    "values": {"type": "array", "items": {"type": "number"}}
                                }
                            }}
                        },
                        "required": ["presentation_id", "title", "chart_type", "categories", "series"]
                    }
                ),
                
                # === ADVANCED SLIDES ===
                types.Tool(
                    name="add_smartart_slide",
                    description="Add a SmartArt diagram slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "smartart_type": {"type": "string", "enum": ["process", "cycle", "hierarchy", "relationship", "matrix"]},
                            "items": {"type": "array", "items": {"type": "string"}}
                        },
                        "required": ["presentation_id", "title", "smartart_type", "items"]
                    }
                ),
                
                types.Tool(
                    name="add_timeline_slide",
                    description="Add a timeline slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "events": {"type": "array", "items": {
                                "type": "object",
                                "properties": {
                                    "date": {"type": "string"},
                                    "event": {"type": "string"}
                                }
                            }}
                        },
                        "required": ["presentation_id", "title", "events"]
                    }
                ),
                
                types.Tool(
                    name="add_comparison_slide",
                    description="Add a comparison slide with two columns",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "left_title": {"type": "string"},
                            "left_points": {"type": "array", "items": {"type": "string"}},
                            "right_title": {"type": "string"},
                            "right_points": {"type": "array", "items": {"type": "string"}}
                        },
                        "required": ["presentation_id", "title"]
                    }
                ),
                
                types.Tool(
                    name="add_quote_slide",
                    description="Add a quote slide with attribution",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "quote": {"type": "string"},
                            "author": {"type": "string"}
                        },
                        "required": ["presentation_id", "quote"]
                    }
                ),
                
                # === ENHANCEMENTS ===
                types.Tool(
                    name="add_image_to_slide",
                    description="Add an image to a specific slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "image_path": {"type": "string"},
                            "position": {"type": "object", "properties": {
                                "left": {"type": "number"}, "top": {"type": "number"},
                                "width": {"type": "number"}, "height": {"type": "number"}
                            }}
                        },
                        "required": ["presentation_id", "slide_index", "image_path"]
                    }
                ),
                
                types.Tool(
                    name="add_qr_code",
                    description="Add a QR code to a slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "url": {"type": "string"},
                            "size": {"type": "number", "default": 1.5}
                        },
                        "required": ["presentation_id", "slide_index", "url"]
                    }
                ),
                
                types.Tool(
                    name="add_watermark",
                    description="Add a watermark to all slides",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "text": {"type": "string"},
                            "opacity": {"type": "number", "default": 0.3}
                        },
                        "required": ["presentation_id", "text"]
                    }
                ),
                
                types.Tool(
                    name="add_slide_notes",
                    description="Add speaker notes to a slide",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "slide_index": {"type": "integer"},
                            "notes": {"type": "string"}
                        },
                        "required": ["presentation_id", "slide_index", "notes"]
                    }
                ),
                
                # === MATERIAL DESIGN ===
                types.Tool(
                    name="apply_material_theme",
                    description="Apply Material Design theme to presentation",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "theme_name": {"type": "string", "enum": ["material_baseline", "google_blue"]},
                            "apply_to_all": {"type": "boolean", "default": True}
                        },
                        "required": ["presentation_id", "theme_name"]
                    }
                ),
                
                types.Tool(
                    name="get_material_color_palette",
                    description="Generate Material Design color palette from seed color",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "seed_color": {"type": "string", "description": "Hex color (e.g., '4CAF50')"}
                        },
                        "required": ["seed_color"]
                    }
                ),
                
                types.Tool(
                    name="check_accessibility",
                    description="Check color accessibility (WCAG compliance)",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "background_color": {"type": "string"},
                            "text_color": {"type": "string"}
                        },
                        "required": ["background_color", "text_color"]
                    }
                ),
                
                # === SAVE & EXPORT ===
                types.Tool(
                    name="save_presentation",
                    description="Save presentation to file",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "presentation_id": {"type": "string"},
                            "file_path": {"type": "string"}
                        },
                        "required": ["presentation_id", "file_path"]
                    }
                )
            ]

        @self.server.call_tool()
        async def handle_call_tool(
            name: str, arguments: dict | None
        ) -> list[types.TextContent | types.ImageContent | types.EmbeddedResource]:
            """Handle tool calls - route to appropriate method"""
            
            if arguments is None:
                arguments = {}
            
            try:
                # Route to appropriate handler
                if name == "create_presentation":
                    return await self.create_presentation(arguments)
                elif name == "convert_markdown_to_pptx":
                    return await self.convert_markdown_to_pptx(arguments)
                elif name == "convert_markdown_file_to_pptx":
                    return await self.convert_markdown_file_to_pptx(arguments)
                elif name == "add_title_slide":
                    return await self.add_title_slide(arguments)
                elif name == "add_content_slide":
                    return await self.add_content_slide(arguments)
                elif name == "add_two_column_slide":
                    return await self.add_two_column_slide(arguments)
                elif name == "add_chart_slide":
                    return await self.add_chart_slide(arguments)
                elif name == "add_smartart_slide":
                    return await self.add_smartart_slide(arguments)
                elif name == "add_timeline_slide":
                    return await self.add_timeline_slide(arguments)
                elif name == "add_comparison_slide":
                    return await self.add_comparison_slide(arguments)
                elif name == "add_quote_slide":
                    return await self.add_quote_slide(arguments)
                elif name == "add_image_to_slide":
                    return await self.add_image_to_slide(arguments)
                elif name == "add_qr_code":
                    return await self.add_qr_code(arguments)
                elif name == "add_watermark":
                    return await self.add_watermark(arguments)
                elif name == "add_slide_notes":
                    return await self.add_slide_notes(arguments)
                elif name == "apply_material_theme":
                    return await self.apply_material_theme(arguments)
                elif name == "get_material_color_palette":
                    return await self.get_material_color_palette(arguments)
                elif name == "check_accessibility":
                    return await self.check_accessibility(arguments)
                elif name == "save_presentation":
                    return await self.save_presentation(arguments)
                else:
                    raise ValueError(f"Unknown tool: {name}")
                    
            except Exception as e:
                return [types.TextContent(
                    type="text",
                    text=f"Error executing {name}: {str(e)}"
                )]

    # ========================================================================
    # PRESENTATION MANAGEMENT
    # ========================================================================

    async def create_presentation(self, args: Dict) -> list[types.TextContent]:
        """Create a new presentation"""
        pres_id = args["presentation_id"]
        template = args.get("template", "corporate")
        aspect_ratio = args.get("aspect_ratio", "16:9")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333 if aspect_ratio == "16:9" else 10)
        prs.slide_height = Inches(7.5 if aspect_ratio == "16:9" else 7.5)
        
        self.presentations[pres_id] = prs
        self.slide_notes[pres_id] = {}
        
        return [types.TextContent(
            type="text",
            text=f"Created presentation '{pres_id}' ({aspect_ratio}, {template} template)"
        )]

    # ========================================================================
    # MARKDOWN CONVERSION
    # ========================================================================

    async def convert_markdown_to_pptx(self, args: Dict) -> list[types.TextContent]:
        """Convert markdown content to PowerPoint"""
        markdown_content = args["markdown_content"]
        output_path = args["output_path"]
        pres_id = args.get("presentation_id", "markdown_conversion")
        
        # Parse markdown
        config, slides = self._parse_markdown(markdown_content)
        
        # Generate presentation
        prs = self._generate_presentation(config, slides)
        
        # Save
        prs.save(output_path)
        
        # Optionally store in server
        if pres_id:
            self.presentations[pres_id] = prs
        
        return [types.TextContent(
            type="text",
            text=f"Converted markdown to PowerPoint: {output_path} ({len(slides)} slides)"
        )]

    async def convert_markdown_file_to_pptx(self, args: Dict) -> list[types.TextContent]:
        """Convert markdown file to PowerPoint"""
        markdown_file = args["markdown_file"]
        output_path = args["output_path"]
        
        with open(markdown_file, 'r', encoding='utf-8') as f:
            markdown_content = f.read()
        
        return await self.convert_markdown_to_pptx({
            "markdown_content": markdown_content,
            "output_path": output_path
        })

    def _parse_markdown(self, markdown_content: str) -> Tuple[PresentationConfig, List[SlideConfig]]:
        """Parse markdown content into config and slides"""
        
        # Extract YAML frontmatter
        config = PresentationConfig(title="Presentation")
        meta_match = re.search(r'^---\s*\n(.*?)\n---\s*\n', markdown_content, re.MULTILINE | re.DOTALL)
        if meta_match:
            try:
                metadata = yaml.safe_load(meta_match.group(1))
                config = PresentationConfig(
                    title=metadata.get('title', 'Presentation'),
                    author=metadata.get('author'),
                    theme=metadata.get('theme', 'default'),
                    aspect_ratio=metadata.get('aspect_ratio', '16:9')
                )
            except yaml.YAMLError:
                pass
        
        # Split into slides by ---
        slides_raw = re.split(r'\n---+\n', markdown_content)
        slides = []
        
        for slide_raw in slides_raw:
            if not slide_raw.strip() or slide_raw.strip().startswith('---'):
                continue
                
            # Parse slide content
            lines = slide_raw.strip().split('\n')
            title = None
            content = []
            
            for line in lines:
                if line.startswith('# '):
                    title = line[2:].strip()
                elif line.startswith('- ') or line.startswith('* '):
                    content.append(line[2:].strip())
                elif line.strip() and not line.startswith('#'):
                    content.append(line.strip())
            
            if title or content:
                slides.append(SlideConfig(
                    type=SlideType.CONTENT if content else SlideType.TITLE,
                    title=title,
                    content=content if content else None
                ))
        
        return config, slides

    def _generate_presentation(self, config: PresentationConfig, slides: List[SlideConfig]) -> Presentation:
        """Generate PowerPoint from parsed config and slides"""
        prs = Presentation()
        
        for slide_config in slides:
            if slide_config.type == SlideType.TITLE or (slide_config.title and not slide_config.content):
                # Title slide
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = slide_config.title or config.title
                if slide.placeholders[1]:
                    slide.placeholders[1].text = slide_config.subtitle or ""
            else:
                # Content slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                if slide_config.title:
                    slide.shapes.title.text = slide_config.title
                if slide_config.content and len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    for i, point in enumerate(slide_config.content):
                        if i == 0:
                            text_frame.text = point
                        else:
                            p = text_frame.add_paragraph()
                            p.text = point
                            p.level = 0
        
        return prs

    # ========================================================================
    # SLIDE CREATION
    # ========================================================================

    async def add_title_slide(self, args: Dict) -> list[types.TextContent]:
        """Add title slide"""
        pres_id = args["presentation_id"]
        title = args["title"]
        subtitle = args.get("subtitle", "")
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        
        return [types.TextContent(type="text", text=f"Added title slide: {title}")]

    async def add_content_slide(self, args: Dict) -> list[types.TextContent]:
        """Add content slide with bullets"""
        pres_id = args["presentation_id"]
        title = args["title"]
        content = args["content"]
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        
        text_frame = slide.placeholders[1].text_frame
        for i, point in enumerate(content):
            if i == 0:
                text_frame.text = point
            else:
                p = text_frame.add_paragraph()
                p.text = point
        
        return [types.TextContent(type="text", text=f"Added content slide: {title}")]

    async def add_two_column_slide(self, args: Dict) -> list[types.TextContent]:
        """Add two-column slide"""
        pres_id = args["presentation_id"]
        title = args["title"]
        left_content = args["left_content"]
        right_content = args["right_content"]
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[3])
        slide.shapes.title.text = title
        
        # Add left column
        left = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.5), Inches(4))
        for item in left_content:
            p = left.text_frame.add_paragraph()
            p.text = f"• {item}"
        
        # Add right column
        right = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(4))
        for item in right_content:
            p = right.text_frame.add_paragraph()
            p.text = f"• {item}"
        
        return [types.TextContent(type="text", text=f"Added two-column slide: {title}")]

    # ========================================================================
    # CHARTS
    # ========================================================================

    async def add_chart_slide(self, args: Dict) -> list[types.TextContent]:
        """Add chart slide"""
        pres_id = args["presentation_id"]
        title = args["title"]
        chart_type = args["chart_type"]
        categories = args["categories"]
        series = args["series"]
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        
        # Prepare chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for s in series:
            chart_data.add_series(s["name"], s["values"])
        
        # Map chart type
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        
        # Add chart
        x, y, cx, cy = Inches(2), Inches(2), Inches(9), Inches(5)
        chart = slide.shapes.add_chart(
            chart_type_map[chart_type], x, y, cx, cy, chart_data
        ).chart
        
        return [types.TextContent(type="text", text=f"Added {chart_type} chart: {title}")]

    # ========================================================================
    # ADVANCED SLIDES
    # ========================================================================

    async def add_smartart_slide(self, args: Dict) -> list[types.TextContent]:
        """Add SmartArt slide"""
        pres_id = args["presentation_id"]
        title = args["title"]
        smartart_type = args["smartart_type"]
        items = args["items"]
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.title.text = title
        
        # Simple SmartArt simulation with shapes
        item_count = len(items)
        width = Inches(10) / item_count
        
        for i, item in enumerate(items):
            left = Inches(1.5 + i * (10 / item_count))
            top = Inches(3)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, Inches(2), Inches(1.5)
            )
            shape.text = item
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(100, 150, 200)
        
        return [types.TextContent(type="text", text=f"Added {smartart_type} SmartArt with {len(items)} items")]

    async def add_timeline_slide(self, args: Dict) -> list[types.TextContent]:
        """Add timeline slide"""
        pres_id = args["presentation_id"]
        title = args["title"]
        events = args["events"]
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.title.text = title
        
        # Draw timeline
        event_count = len(events)
        for i, event in enumerate(events):
            left = Inches(1 + i * (11 / event_count))
            top = Inches(3.5)
            
            # Event marker
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(50, 100, 200)
            
            # Event text
            text_box = slide.shapes.add_textbox(left - Inches(0.5), top + Inches(0.8), Inches(1.5), Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = f"{event['date']}\n{event['event']}"
            text_frame.paragraphs[0].font.size = Pt(10)
        
        return [types.TextContent(type="text", text=f"Added timeline with {len(events)} events")]

    async def add_comparison_slide(self, args: Dict) -> list[types.TextContent]:
        """Add comparison slide"""
        pres_id = args["presentation_id"]
        title = args["title"]
        left_title = args.get("left_title", "Option A")
        left_points = args.get("left_points", [])
        right_title = args.get("right_title", "Option B")
        right_points = args.get("right_points", [])
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.title.text = title
        
        # Left side
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
        tf = left_box.text_frame
        tf.text = left_title
        tf.paragraphs[0].font.bold = True
        for point in left_points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
        
        # Right side
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5))
        tf = right_box.text_frame
        tf.text = right_title
        tf.paragraphs[0].font.bold = True
        for point in right_points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
        
        return [types.TextContent(type="text", text=f"Added comparison slide: {title}")]

    async def add_quote_slide(self, args: Dict) -> list[types.TextContent]:
        """Add quote slide"""
        pres_id = args["presentation_id"]
        quote = args["quote"]
        author = args.get("author", "")
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Quote text
        quote_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
        tf = quote_box.text_frame
        tf.text = f'"{quote}"'
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Author
        if author:
            author_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(0.5))
            af = author_box.text_frame
            af.text = f"— {author}"
            af.paragraphs[0].font.size = Pt(18)
            af.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return [types.TextContent(type="text", text=f"Added quote slide from {author}")]

    # ========================================================================
    # ENHANCEMENTS
    # ========================================================================

    async def add_image_to_slide(self, args: Dict) -> list[types.TextContent]:
        """Add image to slide"""
        pres_id = args["presentation_id"]
        slide_index = args["slide_index"]
        image_path = args["image_path"]
        position = args.get("position", {})
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides[slide_index]
        
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 2))
        width = Inches(position.get("width", 5))
        height = Inches(position.get("height", 3.5))
        
        slide.shapes.add_picture(image_path, left, top, width, height)
        
        return [types.TextContent(type="text", text=f"Added image to slide {slide_index}")]

    async def add_qr_code(self, args: Dict) -> list[types.TextContent]:
        """Add QR code to slide"""
        pres_id = args["presentation_id"]
        slide_index = args["slide_index"]
        url = args["url"]
        size = args.get("size", 1.5)
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        # Generate QR code
        try:
            qr = qrcode.QRCode(version=1, box_size=10, border=4)
            qr.add_data(url)
            qr.make(fit=True)
            img = qr.make_image(fill_color="black", back_color="white")
            
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                img.save(tmp.name)
                tmp_path = tmp.name
            
            # Add to slide
            slide = prs.slides[slide_index]
            slide.shapes.add_picture(tmp_path, Inches(11), Inches(6), Inches(size))
            
            # Cleanup
            Path(tmp_path).unlink()
            
            return [types.TextContent(type="text", text=f"Added QR code to slide {slide_index}")]
        except ImportError:
            return [types.TextContent(type="text", text="QR code generation requires 'qrcode' library")]

    async def add_watermark(self, args: Dict) -> list[types.TextContent]:
        """Add watermark to all slides"""
        pres_id = args["presentation_id"]
        text = args["text"]
        opacity = args.get("opacity", 0.3)
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        for slide in prs.slides:
            textbox = slide.shapes.add_textbox(
                Inches(3), Inches(3.5), Inches(7), Inches(1)
            )
            tf = textbox.text_frame
            tf.text = text
            p = tf.paragraphs[0]
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER
            
            # Rotate (limited support in python-pptx)
            textbox.rotation = 315
        
        return [types.TextContent(type="text", text=f"Added watermark '{text}' to all slides")]

    async def add_slide_notes(self, args: Dict) -> list[types.TextContent]:
        """Add notes to slide"""
        pres_id = args["presentation_id"]
        slide_index = args["slide_index"]
        notes = args["notes"]
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        slide = prs.slides[slide_index]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
        
        if pres_id not in self.slide_notes:
            self.slide_notes[pres_id] = {}
        self.slide_notes[pres_id][slide_index] = notes
        
        return [types.TextContent(type="text", text=f"Added notes to slide {slide_index}")]

    # ========================================================================
    # MATERIAL DESIGN
    # ========================================================================

    async def apply_material_theme(self, args: Dict) -> list[types.TextContent]:
        """Apply Material Design theme"""
        pres_id = args["presentation_id"]
        theme_name = args["theme_name"]
        apply_to_all = args.get("apply_to_all", True)
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        theme = self.material_themes.get(theme_name)
        if not theme:
            return [types.TextContent(type="text", text=f"Theme '{theme_name}' not found")]
        
        # Apply theme colors to all slides
        for slide in prs.slides:
            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(theme.background))
            
            # Update text colors
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(*self._hex_to_rgb(theme.on_background))
        
        return [types.TextContent(
            type="text",
            text=f"Applied Material theme '{theme_name}' to presentation"
        )]

    async def get_material_color_palette(self, args: Dict) -> list[types.TextContent]:
        """Generate color palette from seed color"""
        seed_color = args["seed_color"]
        
        # Generate palette
        palette = self._generate_color_palette(seed_color)
        
        result = {
            "seed_color": seed_color,
            "primary": palette["primary"],
            "secondary": palette["secondary"],
            "accent": palette["accent"],
            "surface": palette["surface"]
        }
        
        return [types.TextContent(
            type="text",
            text=json.dumps(result, indent=2)
        )]

    async def check_accessibility(self, args: Dict) -> list[types.TextContent]:
        """Check color accessibility"""
        bg_color = args["background_color"]
        text_color = args["text_color"]
        
        # Calculate contrast ratio
        contrast = self._calculate_contrast_ratio(bg_color, text_color)
        
        # WCAG levels
        aa_normal = contrast >= 4.5
        aa_large = contrast >= 3.0
        aaa_normal = contrast >= 7.0
        aaa_large = contrast >= 4.5
        
        result = {
            "contrast_ratio": round(contrast, 2),
            "wcag_aa_normal": aa_normal,
            "wcag_aa_large": aa_large,
            "wcag_aaa_normal": aaa_normal,
            "wcag_aaa_large": aaa_large,
            "recommendation": "PASS" if aa_normal else "FAIL"
        }
        
        return [types.TextContent(
            type="text",
            text=json.dumps(result, indent=2)
        )]

    # ========================================================================
    # SAVE & EXPORT
    # ========================================================================

    async def save_presentation(self, args: Dict) -> list[types.TextContent]:
        """Save presentation to file"""
        pres_id = args["presentation_id"]
        file_path = args["file_path"]
        
        prs = self.presentations.get(pres_id)
        if not prs:
            raise ValueError(f"Presentation '{pres_id}' not found")
        
        prs.save(file_path)
        file_size = Path(file_path).stat().st_size
        
        return [types.TextContent(
            type="text",
            text=f"Saved presentation to {file_path} ({file_size:,} bytes)"
        )]

    # ========================================================================
    # UTILITY METHODS
    # ========================================================================

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _generate_color_palette(self, seed_color: str) -> Dict[str, str]:
        """Generate color palette from seed color"""
        rgb = self._hex_to_rgb(seed_color)
        h, l, s = colorsys.rgb_to_hls(rgb[0]/255, rgb[1]/255, rgb[2]/255)
        
        def hls_to_hex(h, l, s):
            r, g, b = colorsys.hls_to_rgb(h, l, s)
            return f"{int(r*255):02x}{int(g*255):02x}{int(b*255):02x}"
        
        return {
            "primary": seed_color,
            "secondary": hls_to_hex((h + 0.083) % 1, l, s),  # +30 degrees
            "accent": hls_to_hex((h + 0.5) % 1, l, s),  # Complementary
            "surface": hls_to_hex(h, 0.95, s * 0.2)
        }

    def _calculate_contrast_ratio(self, color1: str, color2: str) -> float:
        """Calculate WCAG contrast ratio"""
        def relative_luminance(hex_color):
            rgb = self._hex_to_rgb(hex_color)
            rgb_norm = [c / 255.0 for c in rgb]
            rgb_linear = [
                c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
                for c in rgb_norm
            ]
            return 0.2126 * rgb_linear[0] + 0.7152 * rgb_linear[1] + 0.0722 * rgb_linear[2]
        
        l1 = relative_luminance(color1)
        l2 = relative_luminance(color2)
        
        lighter = max(l1, l2)
        darker = min(l1, l2)
        
        return (lighter + 0.05) / (darker + 0.05)

    # ========================================================================
    # SERVER RUNNER
    # ========================================================================

    async def run(self):
        """Run the MCP server"""
        async with stdio_server() as (read_stream, write_stream):
            await self.server.run(
                read_stream,
                write_stream,
                InitializationOptions(
                    server_name="unified-powerpoint-server",
                    server_version="1.0.0",
                    capabilities=self.server.get_capabilities(
                        notification_options=NotificationOptions(),
                        experimental_capabilities={}
                    )
                )
            )


# ============================================================================
# MAIN
# ============================================================================

async def main():
    """Main entry point"""
    server = UnifiedPowerPointServer()
    await server.run()


if __name__ == "__main__":
    asyncio.run(main())
